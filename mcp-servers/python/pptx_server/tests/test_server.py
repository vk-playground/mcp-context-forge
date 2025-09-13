# -*- coding: utf-8 -*-
"""Tests for the PowerPoint MCP Server."""

# Standard
import asyncio
import json
import os
import tempfile

# Third-Party
from pptx import Presentation
from pptx_server.server import (
    add_chart,
    add_shape,
    add_slide,
    add_table,
    add_text_box,
    call_tool,
    create_presentation,
    get_presentation_info,
    list_shapes,
    list_slides,
    save_presentation,
    set_slide_title,
    set_table_cell,
)
import pytest


class TestPresentationBasics:
    """Test basic presentation operations."""

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir

    @pytest.fixture
    def test_pptx_path(self, temp_dir):
        """Return a test PowerPoint file path."""
        return os.path.join(temp_dir, "test_presentation.pptx")

    async def test_create_presentation(self, test_pptx_path):
        """Test creating a new presentation."""
        result = await create_presentation(test_pptx_path, "Test Presentation")

        assert result["message"] == f"Created presentation: {test_pptx_path}"
        assert result["slide_count"] == 1  # Title slide
        assert os.path.exists(test_pptx_path)

        # Verify it's a valid PowerPoint file
        prs = Presentation(test_pptx_path)
        assert len(prs.slides) == 1

    async def test_create_presentation_without_title(self, test_pptx_path):
        """Test creating a presentation without a title."""
        result = await create_presentation(test_pptx_path)

        assert result["message"] == f"Created presentation: {test_pptx_path}"
        assert result["slide_count"] == 0  # No slides added
        assert os.path.exists(test_pptx_path)

    async def test_get_presentation_info(self, test_pptx_path):
        """Test getting presentation information."""
        await create_presentation(test_pptx_path, "Test Presentation")
        result = await get_presentation_info(test_pptx_path)

        assert result["file_path"] == test_pptx_path
        assert result["slide_count"] == 1
        assert result["layout_count"] > 0

    async def test_save_presentation(self, test_pptx_path):
        """Test saving a presentation."""
        await create_presentation(test_pptx_path, "Test Presentation")
        result = await save_presentation(test_pptx_path)

        assert result["message"] == f"Saved presentation: {test_pptx_path}"
        assert os.path.exists(test_pptx_path)


class TestSlideOperations:
    """Test slide management operations."""

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir

    @pytest.fixture
    def test_pptx_path(self, temp_dir):
        """Return a test PowerPoint file path."""
        return os.path.join(temp_dir, "test_slides.pptx")

    @pytest.fixture
    async def presentation_with_slides(self, test_pptx_path):
        """Create a presentation with some slides."""
        await create_presentation(test_pptx_path, "Test Presentation")
        await add_slide(test_pptx_path, 1)  # Content slide
        await add_slide(test_pptx_path, 1)  # Another content slide
        return test_pptx_path

    async def test_add_slide(self, test_pptx_path):
        """Test adding a slide to a presentation."""
        await create_presentation(test_pptx_path)
        result = await add_slide(test_pptx_path, 0)  # Title slide layout

        assert "Added slide at position" in result["message"]
        assert result["slide_index"] == 0

    async def test_list_slides(self, presentation_with_slides):
        """Test listing slides in a presentation."""
        result = await list_slides(presentation_with_slides)

        assert result["total_count"] == 3  # Title + 2 content slides
        assert len(result["slides"]) == 3
        assert all("index" in slide for slide in result["slides"])

    async def test_set_slide_title(self, presentation_with_slides):
        """Test setting slide title."""
        result = await set_slide_title(presentation_with_slides, 0, "New Title")

        assert "Set title for slide 0: New Title" in result["message"]

    async def test_slide_index_validation(self, presentation_with_slides):
        """Test slide index validation."""
        with pytest.raises(ValueError, match="Slide index.*out of range"):
            await set_slide_title(presentation_with_slides, 999, "Invalid")


class TestContentOperations:
    """Test content management operations."""

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir

    @pytest.fixture
    def test_pptx_path(self, temp_dir):
        """Return a test PowerPoint file path."""
        return os.path.join(temp_dir, "test_content.pptx")

    @pytest.fixture
    async def presentation_with_content_slide(self, test_pptx_path):
        """Create a presentation with a content slide."""
        await create_presentation(test_pptx_path)
        await add_slide(test_pptx_path, 1)  # Content slide layout
        return test_pptx_path

    async def test_add_text_box(self, presentation_with_content_slide):
        """Test adding a text box to a slide."""
        result = await add_text_box(
            presentation_with_content_slide,
            slide_index=0,
            text="Test text box",
            left=1.0,
            top=1.0,
            width=4.0,
            height=1.0,
            font_size=18,
            bold=True,
        )

        assert "Added text box to slide 0" in result["message"]
        assert result["text"] == "Test text box"
        assert "shape_index" in result

    async def test_add_shape(self, presentation_with_content_slide):
        """Test adding a shape to a slide."""
        result = await add_shape(
            presentation_with_content_slide,
            slide_index=0,
            shape_type="rectangle",
            left=2.0,
            top=2.0,
            width=3.0,
            height=2.0,
            fill_color="#FF0000",
            line_color="#000000",
        )

        assert "Added rectangle shape to slide 0" in result["message"]
        assert "shape_index" in result

    async def test_invalid_shape_type(self, presentation_with_content_slide):
        """Test adding an invalid shape type."""
        with pytest.raises(ValueError, match="Unknown shape type"):
            await add_shape(
                presentation_with_content_slide,
                slide_index=0,
                shape_type="invalid_shape",
            )


class TestTableOperations:
    """Test table creation and manipulation."""

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir

    @pytest.fixture
    def test_pptx_path(self, temp_dir):
        """Return a test PowerPoint file path."""
        return os.path.join(temp_dir, "test_tables.pptx")

    @pytest.fixture
    async def presentation_with_table(self, test_pptx_path):
        """Create a presentation with a table."""
        await create_presentation(test_pptx_path)
        await add_slide(test_pptx_path, 1)  # Content slide
        await add_table(test_pptx_path, 0, rows=3, cols=4)
        return test_pptx_path

    async def test_add_table(self, test_pptx_path):
        """Test adding a table to a slide."""
        await create_presentation(test_pptx_path)
        await add_slide(test_pptx_path, 1)

        result = await add_table(test_pptx_path, 0, rows=3, cols=4, left=1.0, top=1.0)

        assert "Added 3x4 table to slide 0" in result["message"]
        assert result["rows"] == 3
        assert result["cols"] == 4
        assert "shape_index" in result

    async def test_set_table_cell(self, presentation_with_table):
        """Test setting table cell content."""
        result = await set_table_cell(presentation_with_table, slide_index=0, table_index=0, row=0, col=0, text="Header 1")

        assert "Set cell [0,0] text: Header 1" in result["message"]

    async def test_table_cell_bounds_checking(self, presentation_with_table):
        """Test table cell bounds checking."""
        with pytest.raises(ValueError, match="Row.*out of range"):
            await set_table_cell(
                presentation_with_table,
                slide_index=0,
                table_index=0,
                row=999,
                col=0,
                text="Invalid",
            )


class TestChartOperations:
    """Test chart creation and manipulation."""

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir

    @pytest.fixture
    def test_pptx_path(self, temp_dir):
        """Return a test PowerPoint file path."""
        return os.path.join(temp_dir, "test_charts.pptx")

    @pytest.fixture
    def sample_chart_data(self):
        """Sample chart data for testing."""
        return {
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": [
                {"name": "Revenue", "values": [100, 150, 120, 200]},
                {"name": "Expenses", "values": [80, 90, 85, 95]},
            ],
        }

    async def test_add_chart(self, test_pptx_path, sample_chart_data):
        """Test adding a chart to a slide."""
        await create_presentation(test_pptx_path)
        await add_slide(test_pptx_path, 1)  # Content slide

        result = await add_chart(
            test_pptx_path,
            slide_index=0,
            data=sample_chart_data,
            chart_type="column",
            title="Test Chart",
        )

        assert "Added column chart to slide 0" in result["message"]
        assert result["title"] == "Test Chart"
        assert "shape_index" in result

    async def test_chart_types(self, test_pptx_path, sample_chart_data):
        """Test different chart types."""
        await create_presentation(test_pptx_path)
        await add_slide(test_pptx_path, 1)

        chart_types = ["column", "bar", "line", "pie"]
        for chart_type in chart_types:
            result = await add_chart(test_pptx_path, slide_index=0, data=sample_chart_data, chart_type=chart_type)
            assert f"Added {chart_type} chart" in result["message"]

    async def test_invalid_chart_type(self, test_pptx_path, sample_chart_data):
        """Test invalid chart type handling."""
        await create_presentation(test_pptx_path)
        await add_slide(test_pptx_path, 1)

        with pytest.raises(ValueError, match="Unknown chart type"):
            await add_chart(
                test_pptx_path,
                slide_index=0,
                data=sample_chart_data,
                chart_type="invalid_chart_type",
            )


class TestUtilityOperations:
    """Test utility and information functions."""

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir

    @pytest.fixture
    def test_pptx_path(self, temp_dir):
        """Return a test PowerPoint file path."""
        return os.path.join(temp_dir, "test_utils.pptx")

    @pytest.fixture
    async def presentation_with_shapes(self, test_pptx_path):
        """Create a presentation with various shapes."""
        await create_presentation(test_pptx_path)
        await add_slide(test_pptx_path, 1)  # Content slide
        await add_text_box(test_pptx_path, 0, "Text Box", 1.0, 1.0)
        await add_shape(test_pptx_path, 0, "rectangle", 2.0, 2.0)
        await add_table(test_pptx_path, 0, 2, 3, 3.0, 3.0)
        return test_pptx_path

    async def test_list_shapes(self, presentation_with_shapes):
        """Test listing shapes on a slide."""
        result = await list_shapes(presentation_with_shapes, slide_index=0)

        assert result["total_count"] >= 3  # At least text box, shape, and table
        assert len(result["shapes"]) >= 3

        # Check that shape information is provided
        for shape in result["shapes"]:
            assert "index" in shape
            assert "type" in shape
            assert "left" in shape
            assert "top" in shape
            assert "width" in shape
            assert "height" in shape


class TestToolIntegration:
    """Test MCP tool integration and error handling."""

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir

    @pytest.fixture
    def test_pptx_path(self, temp_dir):
        """Return a test PowerPoint file path."""
        return os.path.join(temp_dir, "test_integration.pptx")

    async def test_call_tool_success(self, test_pptx_path):
        """Test successful tool call through the MCP interface."""
        result = await call_tool("create_presentation", {"file_path": test_pptx_path, "title": "Test"})

        assert len(result) == 1
        assert result[0].type == "text"
        response = json.loads(result[0].text)
        assert response["ok"] is True
        assert "result" in response

    async def test_call_tool_error(self):
        """Test tool call error handling."""
        result = await call_tool("create_presentation", {"file_path": "/invalid/path/test.pptx"})

        assert len(result) == 1
        assert result[0].type == "text"
        response = json.loads(result[0].text)
        assert response["ok"] is False
        assert "error" in response

    async def test_call_tool_unknown(self):
        """Test unknown tool handling."""
        result = await call_tool("unknown_tool", {})

        assert len(result) == 1
        assert result[0].type == "text"
        response = json.loads(result[0].text)
        assert response["ok"] is False
        assert "unknown tool" in response["error"].lower()

    async def test_parameter_validation(self, test_pptx_path):
        """Test parameter validation in tool calls."""
        # Test missing required parameter
        with pytest.raises(TypeError):
            await call_tool("create_presentation", {})

        # Test invalid slide index
        await create_presentation(test_pptx_path)
        result = await call_tool("set_slide_title", {"file_path": test_pptx_path, "slide_index": 999, "title": "Test"})

        response = json.loads(result[0].text)
        assert response["ok"] is False
        assert "out of range" in response["error"]


class TestFileHandling:
    """Test file handling and edge cases."""

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir

    async def test_nonexistent_file_handling(self):
        """Test handling of nonexistent files."""
        result = await call_tool("get_presentation_info", {"file_path": "/nonexistent/path.pptx"})

        response = json.loads(result[0].text)
        # Should create a new presentation if file doesn't exist
        assert response["ok"] is True

    async def test_invalid_image_path(self, temp_dir):
        """Test handling of invalid image paths."""
        pptx_path = os.path.join(temp_dir, "test.pptx")
        await create_presentation(pptx_path)
        await add_slide(pptx_path, 1)

        result = await call_tool("add_image", {"file_path": pptx_path, "slide_index": 0, "image_path": "/nonexistent/image.png"})

        response = json.loads(result[0].text)
        assert response["ok"] is False
        assert "not found" in response["error"].lower()

    async def test_concurrent_operations(self, temp_dir):
        """Test concurrent operations on the same presentation."""
        pptx_path = os.path.join(temp_dir, "concurrent_test.pptx")

        # Create presentation
        await create_presentation(pptx_path)

        # Run multiple operations concurrently
        tasks = [
            add_slide(pptx_path, 1),
            add_slide(pptx_path, 1),
            add_slide(pptx_path, 1),
        ]

        results = await asyncio.gather(*tasks)

        # All operations should succeed
        for result in results:
            assert "Added slide" in result["message"]

        # Verify final state
        info = await get_presentation_info(pptx_path)
        assert info["slide_count"] == 3
