#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint MCP Server Demo Script

This script demonstrates all the capabilities of the PowerPoint MCP Server
by creating a comprehensive presentation with all supported features.
"""

# Standard
import asyncio
import os
from pathlib import Path
import sys

# Add src to path for imports
sys.path.insert(0, str(Path(__file__).parent / "src"))

# Third-Party
from pptx_server.server import (
    add_chart,
    add_shape,
    add_slide,
    add_table,
    add_text_box,
    create_presentation,
    get_presentation_info,
    list_shapes,
    list_slides,
    save_presentation,
    set_slide_content,
    set_slide_title,
    set_table_cell,
)


async def create_demo_presentation():
    """Create a comprehensive demo presentation showcasing all features."""
    print("🎯 Creating PowerPoint MCP Server Demo...")

    demo_file = "examples/demos/pptx_mcp_demo.pptx"

    try:
        # 1. Create presentation
        await create_presentation(demo_file, "PowerPoint MCP Server")
        print("✅ Created presentation with title slide")

        # 2. Overview slide with bullet points
        await add_slide(demo_file, 1)  # Content layout
        await set_slide_title(demo_file, 1, "Comprehensive PowerPoint Automation")
        await set_slide_content(
            demo_file,
            1,
            "Complete slide management (create, delete, reorder)\\n"
            "Advanced text formatting and styling\\n"
            "Dynamic shape creation and editing\\n"
            "Professional table operations\\n"
            "Interactive chart generation\\n"
            "Image handling and positioning",
        )
        print("✅ Added overview slide")

        # 3. Text formatting showcase
        await add_slide(demo_file, 1)
        await set_slide_title(demo_file, 2, "Text Formatting Showcase")

        # Add various formatted text boxes
        await add_text_box(demo_file, 2, "BOLD RED HEADING", 1.0, 2.0, 8.0, 0.8, 28, "#FF0000", True, False)
        await add_text_box(demo_file, 2, "Italic blue subtitle", 1.0, 3.0, 8.0, 0.6, 20, "#0066CC", False, True)
        await add_text_box(demo_file, 2, "Regular black body text for detailed content", 1.0, 3.8, 8.0, 0.6, 16, "#000000", False, False)
        await add_text_box(demo_file, 2, "Bold italic green highlight", 1.0, 4.6, 8.0, 0.6, 18, "#00AA00", True, True)
        print("✅ Added text formatting slide")

        # 4. Shape gallery
        await add_slide(demo_file, 1)
        await set_slide_title(demo_file, 3, "Shape Gallery")

        shapes_data = [
            ("rectangle", "#FF6600", 1.0),
            ("oval", "#6600FF", 2.8),
            ("triangle", "#00FF66", 4.6),
            ("diamond", "#FF0066", 6.4),
            ("star", "#0066FF", 8.2),
        ]

        for i, (shape_type, color, x_pos) in enumerate(shapes_data):
            await add_shape(demo_file, 3, shape_type, x_pos, 2.5, 1.6, 1.8, color, "#000000", 2.0)

        await add_text_box(demo_file, 3, "Geometric shapes with custom colors and borders", 1.0, 4.8, 8.0, 0.5, 14, "#666666", False, False)
        print("✅ Added shape gallery slide")

        # 5. Data table
        await add_slide(demo_file, 1)
        await set_slide_title(demo_file, 4, "Sales Performance Table")

        table_result = await add_table(demo_file, 4, 5, 4, 1.5, 2.0, 7.0, 3.0)
        table_idx = table_result["shape_index"]

        # Table headers
        headers = ["Region", "Q1 Sales", "Q2 Sales", "Growth"]
        for col, header in enumerate(headers):
            await set_table_cell(demo_file, 4, table_idx, 0, col, header)

        # Table data
        data_rows = [
            ["North", "$125K", "$145K", "+16%"],
            ["South", "$98K", "$112K", "+14%"],
            ["East", "$156K", "$189K", "+21%"],
            ["West", "$134K", "$151K", "+13%"],
        ]

        for row_idx, row_data in enumerate(data_rows, 1):
            for col_idx, cell_data in enumerate(row_data):
                await set_table_cell(demo_file, 4, table_idx, row_idx, col_idx, cell_data)

        print("✅ Added data table slide")

        # 6. Multiple charts
        await add_slide(demo_file, 1)
        await set_slide_title(demo_file, 5, "Performance Analytics")

        # Column chart
        chart_data = {
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": [
                {"name": "Revenue", "values": [125, 145, 160, 180]},
                {"name": "Target", "values": [120, 140, 155, 175]},
            ],
        }
        await add_chart(demo_file, 5, chart_data, "column", 0.5, 2.0, 4.5, 3.0, "Quarterly Performance")

        # Pie chart data for second chart
        pie_data = {
            "categories": ["Product A", "Product B", "Product C"],
            "series": [{"name": "Market Share", "values": [45, 30, 25]}],
        }
        await add_chart(demo_file, 5, pie_data, "pie", 5.5, 2.0, 4.0, 3.0, "Market Share")

        print("✅ Added charts slide")

        # 7. Summary slide
        await add_slide(demo_file, 1)
        await set_slide_title(demo_file, 6, "PowerPoint MCP Server Summary")
        await set_slide_content(
            demo_file,
            6,
            "✅ 30+ comprehensive tools for PowerPoint automation\\n"
            "✅ Full slide lifecycle management\\n"
            "✅ Advanced text formatting and styling\\n"
            "✅ Dynamic shapes, tables, and charts\\n"
            "✅ Professional presentation generation\\n"
            "✅ Model Context Protocol integration",
        )

        await add_text_box(demo_file, 6, "Ready for Claude Desktop, VS Code, and any MCP client!", 1.0, 5.5, 8.0, 0.8, 16, "#0066CC", True, False)
        print("✅ Added summary slide")

        # 8. Save and get final stats
        await save_presentation(demo_file)

        # Get presentation statistics
        info = await get_presentation_info(demo_file)
        slides_info = await list_slides(demo_file)

        print(f"\n🎉 DEMO COMPLETE!")
        print(f"📄 Created: {demo_file}")
        print(f"📊 Slides: {info['slide_count']}")
        print(f"💾 Size: {os.path.getsize(demo_file):,} bytes")

        print(f"\n📋 Slide Summary:")
        for slide in slides_info["slides"]:
            shapes_info = await list_shapes(demo_file, slide["index"])
            print(f"   {slide['index']}: {slide['title']} ({shapes_info['total_count']} elements)")

        print(f"\n✨ The demo presentation showcases all major features of the PowerPoint MCP Server!")
        print(f"   Open '{demo_file}' to see the results.")

        return demo_file

    except Exception as e:
        print(f"❌ Error creating demo: {e}")
        # Standard
        import traceback

        traceback.print_exc()
        return None


async def main():
    """Main demo function."""
    demo_file = await create_demo_presentation()

    if demo_file and os.path.exists(demo_file):
        print(f"\n🏆 Demo successful! Open {demo_file} in PowerPoint to see the results.")

        # Verify with python-pptx
        try:
            # Third-Party
            from pptx import Presentation

            prs = Presentation(demo_file)

            print(f"\n🔍 Verification:")
            print(f"   Valid PowerPoint file: ✅")
            print(f"   Total slides: {len(prs.slides)}")

            # Count elements
            total_shapes = sum(len(slide.shapes) for slide in prs.slides)
            print(f"   Total elements: {total_shapes}")

        except Exception as e:
            print(f"⚠️  Verification failed: {e}")
    else:
        print("💥 Demo failed!")
        return 1

    return 0


if __name__ == "__main__":
    exit_code = asyncio.run(main())
    sys.exit(exit_code)
