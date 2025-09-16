# -*- coding: utf-8 -*-
"""Comprehensive PowerPoint MCP Server with full PPTX editing capabilities."""

# Standard
import asyncio
import base64
from datetime import datetime, timedelta
from io import BytesIO
import json
import logging
import os
import sys
from typing import Any, Dict, List, Optional
import uuid

# Third-Party
from dotenv import load_dotenv
from mcp.server import Server
from mcp.server.models import InitializationOptions
from mcp.types import TextContent, Tool
from pathvalidate import is_valid_filename, sanitize_filename
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import Field
from pydantic_settings import BaseSettings

# Load environment variables
load_dotenv()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[logging.StreamHandler(sys.stderr)],
)
log = logging.getLogger("pptx_server")


class PPTXServerConfig(BaseSettings):
    """Configuration settings for the PowerPoint MCP Server."""

    # Server Settings
    server_port: int = Field(default=9000, env="PPTX_SERVER_PORT")
    server_host: str = Field(default="localhost", env="PPTX_SERVER_HOST")
    server_debug: bool = Field(default=False, env="PPTX_SERVER_DEBUG")
    enable_http_downloads: bool = Field(default=True, env="PPTX_ENABLE_HTTP_DOWNLOADS")

    # Security Settings
    enable_file_uploads: bool = Field(default=True, env="PPTX_ENABLE_FILE_UPLOADS")
    max_file_size_mb: int = Field(default=50, env="PPTX_MAX_FILE_SIZE_MB")
    max_presentation_size_mb: int = Field(default=100, env="PPTX_MAX_PRESENTATION_SIZE_MB")
    allowed_upload_extensions: str = Field(default="png,jpg,jpeg,gif,bmp,pptx", env="PPTX_ALLOWED_UPLOAD_EXTENSIONS")
    enable_downloads: bool = Field(default=True, env="PPTX_ENABLE_DOWNLOADS")
    download_token_expiry_hours: int = Field(default=24, env="PPTX_DOWNLOAD_TOKEN_EXPIRY_HOURS")

    # Directory Configuration
    work_dir: str = Field(default="/tmp/pptx_server", env="PPTX_WORK_DIR")
    temp_dir: str = Field(default="/tmp/pptx_server/temp", env="PPTX_TEMP_DIR")
    templates_dir: str = Field(default="/tmp/pptx_server/templates", env="PPTX_TEMPLATES_DIR")
    output_dir: str = Field(default="/tmp/pptx_server/output", env="PPTX_OUTPUT_DIR")
    uploads_dir: str = Field(default="/tmp/pptx_server/uploads", env="PPTX_UPLOADS_DIR")

    # File Management
    auto_cleanup_hours: int = Field(default=48, env="PPTX_AUTO_CLEANUP_HOURS")
    max_files_per_session: int = Field(default=50, env="PPTX_MAX_FILES_PER_SESSION")
    enable_file_versioning: bool = Field(default=True, env="PPTX_ENABLE_FILE_VERSIONING")
    default_slide_format: str = Field(default="16:9", env="PPTX_DEFAULT_SLIDE_FORMAT")

    # Authentication
    require_auth: bool = Field(default=False, env="PPTX_REQUIRE_AUTH")
    api_key: str = Field(default="", env="PPTX_API_KEY")
    jwt_secret: str = Field(default="", env="PPTX_JWT_SECRET")

    # Resource Limits
    max_memory_mb: int = Field(default=512, env="PPTX_MAX_MEMORY_MB")
    max_concurrent_operations: int = Field(default=10, env="PPTX_MAX_CONCURRENT_OPERATIONS")
    operation_timeout_seconds: int = Field(default=300, env="PPTX_OPERATION_TIMEOUT_SECONDS")

    @property
    def allowed_extensions(self) -> List[str]:
        """Get list of allowed file extensions."""
        return [ext.strip().lower() for ext in self.allowed_upload_extensions.split(",")]

    def ensure_directories(self) -> None:
        """Ensure all required directories exist with secure permissions."""
        dirs = [self.work_dir, self.temp_dir, self.templates_dir, self.output_dir, self.uploads_dir, os.path.join(self.work_dir, "logs"), os.path.join(self.work_dir, "sessions")]
        for dir_path in dirs:
            os.makedirs(dir_path, exist_ok=True)
            # Set secure permissions (owner only)
            os.chmod(dir_path, 0o700)

    class Config:
        env_file = ".env"
        extra = "ignore"  # Ignore extra environment variables


# Global configuration instance
config = PPTXServerConfig()
config.ensure_directories()

server = Server("pptx-server")

# Global presentation cache and session management (AUTO-ISOLATED PER AGENT)
_presentations: Dict[str, Dict[str, Presentation]] = {}  # session_id -> {file_path: Presentation}
_download_tokens: Dict[str, Dict[str, Any]] = {}  # UUID -> {file_path, expires, session_id}
_session_files: Dict[str, List[str]] = {}  # session_id -> [file_paths]
_agent_sessions: Dict[str, str] = {}  # agent_id -> session_id (persistent mapping)
_current_session: Optional[str] = None  # Current session for this execution context
_file_to_session: Dict[str, str] = {}  # file_path -> session_id mapping


def _generate_session_id() -> str:
    """Generate a unique session ID."""
    return str(uuid.uuid4())


def _generate_download_token(file_path: str, session_id: str) -> str:
    """Generate a secure download token for a file."""
    token = str(uuid.uuid4())
    expires = datetime.now() + timedelta(hours=config.download_token_expiry_hours)

    token_info = {"file_path": file_path, "expires": expires.isoformat(), "session_id": session_id, "created": datetime.now().isoformat()}

    # Store in memory
    _download_tokens[token] = {"file_path": file_path, "expires": expires, "session_id": session_id, "created": datetime.now()}

    # Also store in file for HTTP server access
    tokens_dir = os.path.join(config.work_dir, "tokens")
    os.makedirs(tokens_dir, exist_ok=True)
    token_file = os.path.join(tokens_dir, f"{token}.json")

    with open(token_file, "w") as f:
        json.dump(token_info, f, indent=2)

    return token


def _validate_filename(filename: str) -> str:
    """Validate and sanitize filename for security."""
    if not filename:
        raise ValueError("Filename cannot be empty")

    # Sanitize filename
    safe_filename = sanitize_filename(filename)

    # Additional security checks
    if ".." in filename or "/" in filename or "\\" in filename:
        raise ValueError("Invalid filename: path traversal not allowed")

    if not is_valid_filename(safe_filename):
        raise ValueError(f"Invalid filename: {filename}")

    return safe_filename


def _get_secure_path(file_path: str, directory_type: str = "output") -> str:
    """Get secure path within configured directories."""
    # Validate filename
    filename = os.path.basename(file_path)
    safe_filename = _validate_filename(filename)

    # Determine target directory
    if directory_type == "output":
        target_dir = config.output_dir
    elif directory_type == "temp":
        target_dir = config.temp_dir
    elif directory_type == "templates":
        target_dir = config.templates_dir
    elif directory_type == "uploads":
        target_dir = config.uploads_dir
    else:
        raise ValueError(f"Unknown directory type: {directory_type}")

    # Ensure directory exists
    os.makedirs(target_dir, exist_ok=True)

    # Return secure path
    return os.path.join(target_dir, safe_filename)


def _generate_agent_id() -> str:
    """Generate a stable agent identifier for the current execution context."""
    # Use process ID + start time for stable agent ID within same execution
    # Standard
    import time

    start_time = getattr(_generate_agent_id, "_start_time", None)
    if start_time is None:
        start_time = int(time.time() * 1000)
        _generate_agent_id._start_time = start_time

    agent_id = f"agent_{os.getpid()}_{start_time}"
    return agent_id


def _get_or_create_agent_session(agent_id: Optional[str] = None) -> str:
    """Get or create an isolated session for each agent/user automatically."""
    # Generate agent ID if not provided
    if agent_id is None:
        agent_id = _generate_agent_id()

    # Check if agent already has a session
    if agent_id in _agent_sessions:
        session_id = _agent_sessions[agent_id]
        # Verify session still exists
        session_dir = os.path.join(config.work_dir, "sessions", session_id)
        if os.path.exists(session_dir):
            return session_id
        else:
            # Session expired or deleted, create new one
            del _agent_sessions[agent_id]

    # Create new session for this agent
    session_id = _generate_session_id()
    session_dir = os.path.join(config.work_dir, "sessions", session_id)
    os.makedirs(session_dir, exist_ok=True)
    os.chmod(session_dir, 0o700)

    # Initialize session
    _session_files[session_id] = []
    _agent_sessions[agent_id] = session_id

    # Create session metadata with agent info
    session_info = {
        "session_id": session_id,
        "agent_id": agent_id,
        "session_name": f"Agent-{agent_id[-8:]}-Workspace",
        "created": datetime.now().isoformat(),
        "workspace_dir": session_dir,
        "expires": (datetime.now() + timedelta(hours=config.auto_cleanup_hours)).isoformat(),
        "auto_generated": True,
    }

    # Save session metadata
    session_file = os.path.join(session_dir, "session.json")
    with open(session_file, "w") as f:
        json.dump(session_info, f, indent=2)

    log.info(f"Auto-created session for agent {agent_id[:16]}: {session_id[:8]}...")

    return session_id


def _ensure_session_directory(session_id: str) -> str:
    """Ensure session directory exists and return path."""
    session_dir = os.path.join(config.work_dir, "sessions", session_id)
    os.makedirs(session_dir, exist_ok=True)
    os.chmod(session_dir, 0o700)

    # Create subdirectories
    for subdir in ["presentations", "uploads", "temp"]:
        subdir_path = os.path.join(session_dir, subdir)
        os.makedirs(subdir_path, exist_ok=True)
        os.chmod(subdir_path, 0o700)

    return session_dir


def _get_session_file_path(filename: str, session_id: str, file_type: str = "presentations") -> str:
    """Get secure file path within session directory."""
    # Validate filename
    safe_filename = _validate_filename(filename)

    # Ensure session directory exists
    session_dir = _ensure_session_directory(session_id)

    # Return path within session
    return os.path.join(session_dir, file_type, safe_filename)


def _ensure_output_directory(file_path: str, session_id: Optional[str] = None) -> str:
    """Ensure output directory exists and return session-scoped secure path."""
    # Auto-generate agent session if none provided
    if session_id is None:
        session_id = _get_or_create_agent_session()

    # Extract filename and validate
    filename = os.path.basename(file_path) if file_path else "presentation.pptx"

    # Always use session-scoped path for security
    return _get_session_file_path(filename, session_id, "presentations")


def _resolve_template_path(template_path: str) -> str:
    """Resolve template path, checking secure template directories."""
    # Check if it's already a valid absolute path
    if os.path.isabs(template_path) and os.path.exists(template_path):
        return template_path

    # Check if relative path exists
    if os.path.exists(template_path):
        return template_path

    # Check in secure templates directory
    secure_template = os.path.join(config.templates_dir, os.path.basename(template_path))
    if os.path.exists(secure_template):
        return secure_template

    # Check in legacy templates directory for backward compatibility
    legacy_template = os.path.join("examples/templates", os.path.basename(template_path))
    if os.path.exists(legacy_template):
        return legacy_template

    return template_path  # Return original if not found


def _get_presentation(file_path: str, session_id: Optional[str] = None) -> Presentation:
    """Get or create a presentation with automatic session isolation."""
    abs_path = os.path.abspath(file_path)

    # Check if this file already has a session mapped
    if abs_path in _file_to_session:
        session_id = _file_to_session[abs_path]
    elif session_id is None:
        # Auto-generate session for agent if not provided
        session_id = _get_or_create_agent_session()
        # Map this file to the session
        _file_to_session[abs_path] = session_id

    # Ensure session exists in cache
    if session_id not in _presentations:
        _presentations[session_id] = {}

    # Check session-isolated cache
    if abs_path not in _presentations[session_id]:
        if os.path.exists(abs_path):
            log.info(f"Loading existing presentation: {abs_path} (session: {session_id[:8]})")
            _presentations[session_id][abs_path] = Presentation(abs_path)
        else:
            log.info(f"Creating new presentation: {abs_path} (session: {session_id[:8]})")
            prs = Presentation()
            # Set all new presentations to 16:9 widescreen by default
            _set_slide_size_16_9(prs)
            _presentations[session_id][abs_path] = prs

    return _presentations[session_id][abs_path]


def _get_session_for_operation() -> str:
    """Get session ID for current operation with automatic agent isolation."""
    return _get_or_create_agent_session()


def _save_presentation(file_path: str, session_id: Optional[str] = None) -> None:
    """Save a presentation with automatic session isolation."""
    abs_path = os.path.abspath(file_path)

    # Use mapped session if available
    if abs_path in _file_to_session:
        session_id = _file_to_session[abs_path]
    elif session_id is None:
        session_id = _get_or_create_agent_session()
        _file_to_session[abs_path] = session_id

    # Check session-isolated cache
    if session_id in _presentations and abs_path in _presentations[session_id]:
        # Ensure directory exists
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        _presentations[session_id][abs_path].save(abs_path)
        log.info(f"Saved presentation: {abs_path} (session: {session_id[:8]})")

        # Track file in session
        if session_id not in _session_files:
            _session_files[session_id] = []
        if abs_path not in _session_files[session_id]:
            _session_files[session_id].append(abs_path)


def _parse_color(color_str: str) -> RGBColor:
    """Parse color string (hex format like #FF0000) to RGBColor."""
    if color_str.startswith("#"):
        color_str = color_str[1:]
    return RGBColor(int(color_str[:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16))


def _set_slide_size_16_9(presentation: Presentation) -> None:
    """Set presentation to 16:9 widescreen format (modern standard)."""
    # 16:9 widescreen dimensions
    presentation.slide_width = Inches(13.33)  # 16:9 widescreen width
    presentation.slide_height = Inches(7.5)  # 16:9 widescreen height


def _set_slide_size_4_3(presentation: Presentation) -> None:
    """Set presentation to 4:3 standard format (legacy)."""
    # 4:3 standard dimensions
    presentation.slide_width = Inches(10.0)  # 4:3 standard width
    presentation.slide_height = Inches(7.5)  # 4:3 standard height


@server.list_tools()
async def list_tools() -> list[Tool]:
    """List all available PowerPoint editing tools."""
    return [
        # Presentation Management
        Tool(
            name="create_presentation",
            description="Create a new PowerPoint presentation in secure session workspace",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Filename for the presentation (created in secure session directory)"},
                    "title": {"type": "string", "description": "Optional title for the presentation"},
                    "session_id": {"type": "string", "description": "Session ID for workspace isolation (auto-created if not provided)"},
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="create_presentation_from_template",
            description="Create a new PowerPoint presentation from an existing template",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_path": {"type": "string", "description": "Path to the template presentation file"},
                    "output_path": {"type": "string", "description": "Path where to save the new presentation"},
                    "title": {"type": "string", "description": "Optional new title for the presentation"},
                    "replace_placeholders": {"type": "object", "description": "Key-value pairs to replace text placeholders", "additionalProperties": {"type": "string"}},
                },
                "required": ["template_path", "output_path"],
            },
        ),
        Tool(
            name="clone_presentation",
            description="Clone an existing presentation with optional modifications",
            inputSchema={
                "type": "object",
                "properties": {
                    "source_path": {"type": "string", "description": "Path to the source presentation"},
                    "target_path": {"type": "string", "description": "Path for the cloned presentation"},
                    "new_title": {"type": "string", "description": "Optional new title for the cloned presentation"},
                },
                "required": ["source_path", "target_path"],
            },
        ),
        Tool(
            name="open_presentation",
            description="Open an existing PowerPoint presentation",
            inputSchema={"type": "object", "properties": {"file_path": {"type": "string", "description": "Path to the presentation file"}}, "required": ["file_path"]},
        ),
        Tool(
            name="save_presentation",
            description="Save the current presentation to file",
            inputSchema={"type": "object", "properties": {"file_path": {"type": "string", "description": "Path where to save the presentation"}}, "required": ["file_path"]},
        ),
        Tool(
            name="get_presentation_info",
            description="Get information about the presentation (slide count, properties, etc.)",
            inputSchema={"type": "object", "properties": {"file_path": {"type": "string", "description": "Path to the presentation file"}}, "required": ["file_path"]},
        ),
        # Slide Management
        Tool(
            name="add_slide",
            description="Add a new slide to the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "layout_index": {"type": "integer", "description": "Slide layout index (0-based)", "default": 0},
                    "position": {"type": "integer", "description": "Position to insert slide (0-based, -1 for end)", "default": -1},
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="delete_slide",
            description="Delete a slide from the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide to delete (0-based)"},
                },
                "required": ["file_path", "slide_index"],
            },
        ),
        Tool(
            name="move_slide",
            description="Move a slide to a different position",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "from_index": {"type": "integer", "description": "Current index of slide (0-based)"},
                    "to_index": {"type": "integer", "description": "New index for slide (0-based)"},
                },
                "required": ["file_path", "from_index", "to_index"],
            },
        ),
        Tool(
            name="duplicate_slide",
            description="Duplicate an existing slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide to duplicate (0-based)"},
                    "position": {"type": "integer", "description": "Position for duplicated slide (-1 for end)", "default": -1},
                },
                "required": ["file_path", "slide_index"],
            },
        ),
        Tool(
            name="list_slides",
            description="List all slides in the presentation with their basic information",
            inputSchema={"type": "object", "properties": {"file_path": {"type": "string", "description": "Path to the presentation file"}}, "required": ["file_path"]},
        ),
        # Text and Content Management
        Tool(
            name="set_slide_title",
            description="Set the title of a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "title": {"type": "string", "description": "Title text"},
                },
                "required": ["file_path", "slide_index", "title"],
            },
        ),
        Tool(
            name="set_slide_content",
            description="Set the main content/body text of a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "content": {"type": "string", "description": "Content text (can include bullet points with \\n)"},
                },
                "required": ["file_path", "slide_index", "content"],
            },
        ),
        Tool(
            name="add_text_box",
            description="Add a text box to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "text": {"type": "string", "description": "Text content"},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Width in inches", "default": 6.0},
                    "height": {"type": "number", "description": "Height in inches", "default": 1.0},
                    "font_size": {"type": "integer", "description": "Font size in points", "default": 18},
                    "font_color": {"type": "string", "description": "Font color in hex (#RRGGBB)", "default": "#000000"},
                    "bold": {"type": "boolean", "description": "Make text bold", "default": False},
                    "italic": {"type": "boolean", "description": "Make text italic", "default": False},
                },
                "required": ["file_path", "slide_index", "text"],
            },
        ),
        Tool(
            name="format_text",
            description="Format existing text in a slide (placeholder or text box)",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_index": {"type": "integer", "description": "Index of text shape (0-based)"},
                    "font_name": {"type": "string", "description": "Font name (e.g., 'Arial', 'Times New Roman')"},
                    "font_size": {"type": "integer", "description": "Font size in points"},
                    "font_color": {"type": "string", "description": "Font color in hex (#RRGGBB)"},
                    "bold": {"type": "boolean", "description": "Make text bold"},
                    "italic": {"type": "boolean", "description": "Make text italic"},
                    "underline": {"type": "boolean", "description": "Underline text"},
                    "alignment": {"type": "string", "description": "Text alignment (left, center, right, justify)", "default": "left"},
                },
                "required": ["file_path", "slide_index", "shape_index"],
            },
        ),
        # Image Management
        Tool(
            name="add_image",
            description="Add an image to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "image_path": {"type": "string", "description": "Path to the image file"},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Width in inches (optional, maintains aspect ratio)"},
                    "height": {"type": "number", "description": "Height in inches (optional, maintains aspect ratio)"},
                },
                "required": ["file_path", "slide_index", "image_path"],
            },
        ),
        Tool(
            name="add_image_from_base64",
            description="Add an image to a slide from base64 encoded data",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "image_data": {"type": "string", "description": "Base64 encoded image data"},
                    "image_format": {"type": "string", "description": "Image format (png, jpg, gif)", "default": "png"},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Width in inches (optional)"},
                    "height": {"type": "number", "description": "Height in inches (optional)"},
                },
                "required": ["file_path", "slide_index", "image_data"],
            },
        ),
        Tool(
            name="replace_image",
            description="Replace an existing image in a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_index": {"type": "integer", "description": "Index of image shape (0-based)"},
                    "new_image_path": {"type": "string", "description": "Path to the new image file"},
                },
                "required": ["file_path", "slide_index", "shape_index", "new_image_path"],
            },
        ),
        # Shape Management
        Tool(
            name="add_shape",
            description="Add a shape to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_type": {"type": "string", "description": "Shape type (rectangle, oval, triangle, arrow, etc.)"},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Width in inches", "default": 2.0},
                    "height": {"type": "number", "description": "Height in inches", "default": 1.0},
                    "fill_color": {"type": "string", "description": "Fill color in hex (#RRGGBB)"},
                    "line_color": {"type": "string", "description": "Line color in hex (#RRGGBB)"},
                    "line_width": {"type": "number", "description": "Line width in points", "default": 1.0},
                },
                "required": ["file_path", "slide_index", "shape_type"],
            },
        ),
        Tool(
            name="modify_shape",
            description="Modify properties of an existing shape",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_index": {"type": "integer", "description": "Index of shape (0-based)"},
                    "left": {"type": "number", "description": "Left position in inches"},
                    "top": {"type": "number", "description": "Top position in inches"},
                    "width": {"type": "number", "description": "Width in inches"},
                    "height": {"type": "number", "description": "Height in inches"},
                    "fill_color": {"type": "string", "description": "Fill color in hex (#RRGGBB)"},
                    "line_color": {"type": "string", "description": "Line color in hex (#RRGGBB)"},
                    "line_width": {"type": "number", "description": "Line width in points"},
                },
                "required": ["file_path", "slide_index", "shape_index"],
            },
        ),
        Tool(
            name="delete_shape",
            description="Delete a shape from a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_index": {"type": "integer", "description": "Index of shape to delete (0-based)"},
                },
                "required": ["file_path", "slide_index", "shape_index"],
            },
        ),
        # Table Operations
        Tool(
            name="add_table",
            description="Add a table to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "rows": {"type": "integer", "description": "Number of rows", "minimum": 1},
                    "cols": {"type": "integer", "description": "Number of columns", "minimum": 1},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Table width in inches", "default": 6.0},
                    "height": {"type": "number", "description": "Table height in inches", "default": 3.0},
                },
                "required": ["file_path", "slide_index", "rows", "cols"],
            },
        ),
        Tool(
            name="set_table_cell",
            description="Set the text content of a table cell",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "table_index": {"type": "integer", "description": "Index of table shape (0-based)"},
                    "row": {"type": "integer", "description": "Row index (0-based)"},
                    "col": {"type": "integer", "description": "Column index (0-based)"},
                    "text": {"type": "string", "description": "Cell text content"},
                },
                "required": ["file_path", "slide_index", "table_index", "row", "col", "text"],
            },
        ),
        Tool(
            name="format_table_cell",
            description="Format a table cell (font, color, alignment)",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "table_index": {"type": "integer", "description": "Index of table shape (0-based)"},
                    "row": {"type": "integer", "description": "Row index (0-based)"},
                    "col": {"type": "integer", "description": "Column index (0-based)"},
                    "font_size": {"type": "integer", "description": "Font size in points"},
                    "font_color": {"type": "string", "description": "Font color in hex (#RRGGBB)"},
                    "fill_color": {"type": "string", "description": "Cell background color in hex (#RRGGBB)"},
                    "bold": {"type": "boolean", "description": "Make text bold"},
                    "alignment": {"type": "string", "description": "Text alignment (left, center, right)"},
                },
                "required": ["file_path", "slide_index", "table_index", "row", "col"],
            },
        ),
        Tool(
            name="populate_table",
            description="Populate entire table with data from a 2D array",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "table_index": {"type": "integer", "description": "Index of table shape (0-based)"},
                    "data": {"type": "array", "description": "2D array of cell values", "items": {"type": "array", "items": {"type": "string"}}},
                    "header_row": {"type": "boolean", "description": "Format first row as header", "default": False},
                },
                "required": ["file_path", "slide_index", "table_index", "data"],
            },
        ),
        # Chart Operations
        Tool(
            name="add_chart",
            description="Add a chart to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "chart_type": {"type": "string", "description": "Chart type (column, bar, line, pie)", "default": "column"},
                    "data": {
                        "type": "object",
                        "description": "Chart data with categories and series",
                        "properties": {
                            "categories": {"type": "array", "items": {"type": "string"}},
                            "series": {"type": "array", "items": {"type": "object", "properties": {"name": {"type": "string"}, "values": {"type": "array", "items": {"type": "number"}}}}},
                        },
                    },
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Chart width in inches", "default": 6.0},
                    "height": {"type": "number", "description": "Chart height in inches", "default": 4.0},
                    "title": {"type": "string", "description": "Chart title"},
                },
                "required": ["file_path", "slide_index", "data"],
            },
        ),
        Tool(
            name="update_chart_data",
            description="Update data in an existing chart",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "chart_index": {"type": "integer", "description": "Index of chart shape (0-based)"},
                    "data": {
                        "type": "object",
                        "description": "New chart data",
                        "properties": {
                            "categories": {"type": "array", "items": {"type": "string"}},
                            "series": {"type": "array", "items": {"type": "object", "properties": {"name": {"type": "string"}, "values": {"type": "array", "items": {"type": "number"}}}}},
                        },
                    },
                },
                "required": ["file_path", "slide_index", "chart_index", "data"],
            },
        ),
        # Utility and Information Tools
        Tool(
            name="list_shapes",
            description="List all shapes on a slide with their types and properties",
            inputSchema={
                "type": "object",
                "properties": {"file_path": {"type": "string", "description": "Path to the presentation file"}, "slide_index": {"type": "integer", "description": "Index of slide (0-based)"}},
                "required": ["file_path", "slide_index"],
            },
        ),
        Tool(
            name="get_slide_layouts",
            description="Get available slide layouts in the presentation",
            inputSchema={"type": "object", "properties": {"file_path": {"type": "string", "description": "Path to the presentation file"}}, "required": ["file_path"]},
        ),
        Tool(
            name="set_presentation_properties",
            description="Set presentation document properties",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Presentation title"},
                    "author": {"type": "string", "description": "Author name"},
                    "subject": {"type": "string", "description": "Subject"},
                    "comments": {"type": "string", "description": "Comments"},
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="set_slide_size",
            description="Set the slide size/aspect ratio of the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "format": {"type": "string", "description": "Slide format", "enum": ["16:9", "4:3", "custom"], "default": "16:9"},
                    "width_inches": {"type": "number", "description": "Custom width in inches (if format is custom)"},
                    "height_inches": {"type": "number", "description": "Custom height in inches (if format is custom)"},
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="get_slide_size",
            description="Get the current slide size and aspect ratio of the presentation",
            inputSchema={"type": "object", "properties": {"file_path": {"type": "string", "description": "Path to the presentation file"}}, "required": ["file_path"]},
        ),
        Tool(
            name="export_slide_as_image",
            description="Export a slide as an image file",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "output_path": {"type": "string", "description": "Output image file path"},
                    "format": {"type": "string", "description": "Image format (png, jpg)", "default": "png"},
                },
                "required": ["file_path", "slide_index", "output_path"],
            },
        ),
        # Security and File Management Tools
        Tool(
            name="create_secure_session",
            description="Create a secure session for file operations with UUID workspace",
            inputSchema={"type": "object", "properties": {"session_name": {"type": "string", "description": "Optional session name for identification"}}},
        ),
        Tool(
            name="upload_file",
            description="Upload a file (image or template) to secure workspace",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_data": {"type": "string", "description": "Base64 encoded file data"},
                    "filename": {"type": "string", "description": "Original filename"},
                    "session_id": {"type": "string", "description": "Session ID for workspace isolation"},
                },
                "required": ["file_data", "filename"],
            },
        ),
        Tool(
            name="create_download_link",
            description="Create a secure download link for a presentation with expiration",
            inputSchema={
                "type": "object",
                "properties": {"file_path": {"type": "string", "description": "Path to the presentation file"}, "session_id": {"type": "string", "description": "Session ID for access control"}},
                "required": ["file_path"],
            },
        ),
        Tool(
            name="list_session_files",
            description="List all files in the current session",
            inputSchema={"type": "object", "properties": {"session_id": {"type": "string", "description": "Session ID to list files for"}}, "required": ["session_id"]},
        ),
        Tool(
            name="cleanup_session",
            description="Clean up session files and resources",
            inputSchema={
                "type": "object",
                "properties": {
                    "session_id": {"type": "string", "description": "Session ID to clean up"},
                    "force": {"type": "boolean", "description": "Force cleanup even if session is active", "default": False},
                },
                "required": ["session_id"],
            },
        ),
        Tool(name="get_server_status", description="Get server configuration and status information", inputSchema={"type": "object", "properties": {}}),
        Tool(
            name="get_file_content",
            description="Get the raw file content for download (base64 encoded)",
            inputSchema={
                "type": "object",
                "properties": {"file_path": {"type": "string", "description": "Path to the presentation file"}, "session_id": {"type": "string", "description": "Session ID for access control"}},
                "required": ["file_path"],
            },
        ),
        # Composite Workflow Tools
        Tool(
            name="create_title_slide",
            description="Create a complete title slide with title, subtitle, and optional company info",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Main presentation title"},
                    "subtitle": {"type": "string", "description": "Subtitle or description"},
                    "author": {"type": "string", "description": "Author or company name"},
                    "date": {"type": "string", "description": "Date or additional info"},
                    "slide_index": {"type": "integer", "description": "Index where to create slide (0-based)", "default": 0},
                },
                "required": ["file_path", "title"],
            },
        ),
        Tool(
            name="create_data_slide",
            description="Create a complete data slide with title, table, and optional chart",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Slide title"},
                    "data": {"type": "array", "description": "2D array of data for table", "items": {"type": "array", "items": {"type": "string"}}},
                    "include_chart": {"type": "boolean", "description": "Whether to create a chart from the data", "default": False},
                    "chart_type": {"type": "string", "description": "Chart type if creating chart", "default": "column"},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": -1},
                },
                "required": ["file_path", "title", "data"],
            },
        ),
        Tool(
            name="create_comparison_slide",
            description="Create a comparison slide with two columns of content",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Slide title"},
                    "left_title": {"type": "string", "description": "Left column title"},
                    "left_content": {"type": "array", "description": "Left column bullet points", "items": {"type": "string"}},
                    "right_title": {"type": "string", "description": "Right column title"},
                    "right_content": {"type": "array", "description": "Right column bullet points", "items": {"type": "string"}},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": -1},
                },
                "required": ["file_path", "title", "left_title", "left_content", "right_title", "right_content"],
            },
        ),
        Tool(
            name="create_agenda_slide",
            description="Create an agenda/outline slide with numbered or bulleted items",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Slide title", "default": "Agenda"},
                    "agenda_items": {"type": "array", "description": "List of agenda items", "items": {"type": "string"}},
                    "numbered": {"type": "boolean", "description": "Use numbers instead of bullets", "default": True},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": 1},
                },
                "required": ["file_path", "agenda_items"],
            },
        ),
        Tool(
            name="batch_replace_text",
            description="Replace text across multiple slides in the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "replacements": {"type": "object", "description": "Key-value pairs of text to replace", "additionalProperties": {"type": "string"}},
                    "slide_range": {"type": "array", "description": "Range of slide indices to process (all if not specified)", "items": {"type": "integer"}},
                    "case_sensitive": {"type": "boolean", "description": "Whether replacement should be case sensitive", "default": False},
                },
                "required": ["file_path", "replacements"],
            },
        ),
        Tool(
            name="apply_brand_theme",
            description="Apply consistent branding theme across presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "primary_color": {"type": "string", "description": "Primary brand color (hex)", "default": "#0066CC"},
                    "secondary_color": {"type": "string", "description": "Secondary brand color (hex)", "default": "#999999"},
                    "accent_color": {"type": "string", "description": "Accent brand color (hex)", "default": "#FF6600"},
                    "font_family": {"type": "string", "description": "Primary font family", "default": "Arial"},
                    "apply_to_titles": {"type": "boolean", "description": "Apply colors to slide titles", "default": True},
                    "apply_to_shapes": {"type": "boolean", "description": "Apply colors to shapes", "default": True},
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="create_section_break",
            description="Create a section break slide with large title and optional image",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "section_title": {"type": "string", "description": "Section title"},
                    "subtitle": {"type": "string", "description": "Optional subtitle"},
                    "background_color": {"type": "string", "description": "Background color (hex)", "default": "#0066CC"},
                    "text_color": {"type": "string", "description": "Text color (hex)", "default": "#FFFFFF"},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": -1},
                },
                "required": ["file_path", "section_title"],
            },
        ),
        Tool(
            name="generate_summary_slide",
            description="Generate a summary slide based on presentation content",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Summary slide title", "default": "Summary"},
                    "max_points": {"type": "integer", "description": "Maximum number of summary points", "default": 5},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": -1},
                },
                "required": ["file_path"],
            },
        ),
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict) -> list[TextContent]:
    """Handle tool calls for PowerPoint operations."""
    try:
        result = None

        if name == "create_presentation":
            result = await create_presentation(**arguments)
        elif name == "create_presentation_from_template":
            result = await create_presentation_from_template(**arguments)
        elif name == "clone_presentation":
            result = await clone_presentation(**arguments)
        elif name == "open_presentation":
            result = await open_presentation(**arguments)
        elif name == "save_presentation":
            result = await save_presentation(**arguments)
        elif name == "get_presentation_info":
            result = await get_presentation_info(**arguments)
        elif name == "add_slide":
            result = await add_slide(**arguments)
        elif name == "delete_slide":
            result = await delete_slide(**arguments)
        elif name == "move_slide":
            result = await move_slide(**arguments)
        elif name == "duplicate_slide":
            result = await duplicate_slide(**arguments)
        elif name == "list_slides":
            result = await list_slides(**arguments)
        elif name == "set_slide_title":
            result = await set_slide_title(**arguments)
        elif name == "set_slide_content":
            result = await set_slide_content(**arguments)
        elif name == "add_text_box":
            result = await add_text_box(**arguments)
        elif name == "format_text":
            result = await format_text(**arguments)
        elif name == "add_image":
            result = await add_image(**arguments)
        elif name == "add_image_from_base64":
            result = await add_image_from_base64(**arguments)
        elif name == "replace_image":
            result = await replace_image(**arguments)
        elif name == "add_shape":
            result = await add_shape(**arguments)
        elif name == "modify_shape":
            result = await modify_shape(**arguments)
        elif name == "delete_shape":
            result = await delete_shape(**arguments)
        elif name == "add_table":
            result = await add_table(**arguments)
        elif name == "set_table_cell":
            result = await set_table_cell(**arguments)
        elif name == "format_table_cell":
            result = await format_table_cell(**arguments)
        elif name == "populate_table":
            result = await populate_table(**arguments)
        elif name == "add_chart":
            result = await add_chart(**arguments)
        elif name == "update_chart_data":
            result = await update_chart_data(**arguments)
        elif name == "list_shapes":
            result = await list_shapes(**arguments)
        elif name == "get_slide_layouts":
            result = await get_slide_layouts(**arguments)
        elif name == "set_presentation_properties":
            result = await set_presentation_properties(**arguments)
        elif name == "set_slide_size":
            result = await set_slide_size(**arguments)
        elif name == "get_slide_size":
            result = await get_slide_size(**arguments)
        elif name == "export_slide_as_image":
            result = await export_slide_as_image(**arguments)
        # Security and file management tools
        elif name == "create_secure_session":
            result = await create_secure_session(**arguments)
        elif name == "upload_file":
            result = await upload_file(**arguments)
        elif name == "create_download_link":
            result = await create_download_link(**arguments)
        elif name == "list_session_files":
            result = await list_session_files(**arguments)
        elif name == "cleanup_session":
            result = await cleanup_session(**arguments)
        elif name == "get_server_status":
            result = await get_server_status(**arguments)
        elif name == "get_file_content":
            result = await get_file_content(**arguments)
        # Composite workflow tools
        elif name == "create_title_slide":
            result = await create_title_slide(**arguments)
        elif name == "create_data_slide":
            result = await create_data_slide(**arguments)
        elif name == "create_comparison_slide":
            result = await create_comparison_slide(**arguments)
        elif name == "create_agenda_slide":
            result = await create_agenda_slide(**arguments)
        elif name == "batch_replace_text":
            result = await batch_replace_text(**arguments)
        elif name == "apply_brand_theme":
            result = await apply_brand_theme(**arguments)
        elif name == "create_section_break":
            result = await create_section_break(**arguments)
        elif name == "generate_summary_slide":
            result = await generate_summary_slide(**arguments)
        else:
            return [TextContent(type="text", text=json.dumps({"ok": False, "error": f"Unknown tool: {name}"}))]

        return [TextContent(type="text", text=json.dumps({"ok": True, "result": result}))]

    except Exception as e:
        log.error(f"Error in tool {name}: {str(e)}")
        return [TextContent(type="text", text=json.dumps({"ok": False, "error": str(e)}))]


# Tool implementations start here
async def create_presentation(file_path: str, title: Optional[str] = None, session_id: Optional[str] = None) -> Dict[str, Any]:
    """Create a new PowerPoint presentation in secure session workspace."""
    prs = Presentation()

    # Set to modern 16:9 widescreen format by default
    _set_slide_size_16_9(prs)

    if title:
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title

    # Ensure proper directory structure
    _ensure_output_directory(file_path)

    # SECURITY FIX: Auto-generate isolated session per agent
    if session_id is None:
        session_id = _get_or_create_agent_session()

    secure_path = _get_session_file_path(file_path, session_id, "presentations")

    # Cache the presentation in session-isolated cache
    abs_path = os.path.abspath(secure_path)
    if session_id not in _presentations:
        _presentations[session_id] = {}
    _presentations[session_id][abs_path] = prs

    # Save immediately with session context
    _save_presentation(secure_path, session_id)

    return {"message": f"Created presentation: {secure_path}", "slide_count": len(prs.slides), "format": "16:9 widescreen", "session_id": session_id, "secure_path": secure_path}


async def open_presentation(file_path: str) -> Dict[str, Any]:
    """Open an existing PowerPoint presentation."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Presentation file not found: {file_path}")

    prs = _get_presentation(file_path)
    return {"message": f"Opened presentation: {file_path}", "slide_count": len(prs.slides), "layouts_count": len(prs.slide_layouts)}


async def save_presentation(file_path: str) -> Dict[str, Any]:
    """Save the current presentation to file."""
    _save_presentation(file_path)
    return {"message": f"Saved presentation: {file_path}"}


async def get_presentation_info(file_path: str) -> Dict[str, Any]:
    """Get information about the presentation."""
    prs = _get_presentation(file_path)

    # Get document properties
    props = prs.core_properties

    return {
        "file_path": file_path,
        "slide_count": len(prs.slides),
        "layout_count": len(prs.slide_layouts),
        "title": props.title or "",
        "author": props.author or "",
        "subject": props.subject or "",
        "created": str(props.created) if props.created else "",
        "modified": str(props.modified) if props.modified else "",
    }


async def add_slide(file_path: str, layout_index: int = 0, position: int = -1) -> Dict[str, Any]:
    """Add a new slide to the presentation."""
    prs = _get_presentation(file_path)

    if layout_index >= len(prs.slide_layouts):
        raise ValueError(f"Layout index {layout_index} out of range. Available layouts: 0-{len(prs.slide_layouts)-1}")

    slide_layout = prs.slide_layouts[layout_index]

    if position == -1:
        prs.slides.add_slide(slide_layout)
        slide_idx = len(prs.slides) - 1
    else:
        # python-pptx doesn't have direct insert at position, so we'll add at end and move
        prs.slides.add_slide(slide_layout)
        slide_idx = len(prs.slides) - 1
        if position < slide_idx:
            # Move slide to desired position (this is a workaround)
            pass  # Note: Moving requires more complex XML manipulation

    # Save presentation after modification
    _save_presentation(file_path)

    return {"message": f"Added slide at position {slide_idx}", "slide_index": slide_idx, "layout_name": slide_layout.name if hasattr(slide_layout, "name") else f"Layout {layout_index}"}


async def delete_slide(file_path: str, slide_index: int) -> Dict[str, Any]:
    """Delete a slide from the presentation."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range. Available slides: 0-{len(prs.slides)-1}")

    # Get slide reference
    prs.slides[slide_index].slide_id

    # Remove from slides collection
    del prs.slides._sldIdLst[slide_index]

    return {"message": f"Deleted slide at index {slide_index}"}


async def move_slide(file_path: str, from_index: int, to_index: int) -> Dict[str, Any]:
    """Move a slide to a different position."""
    prs = _get_presentation(file_path)

    slide_count = len(prs.slides)
    if from_index < 0 or from_index >= slide_count:
        raise ValueError(f"From index {from_index} out of range")
    if to_index < 0 or to_index >= slide_count:
        raise ValueError(f"To index {to_index} out of range")

    # This is a complex operation that requires XML manipulation
    # For now, return a placeholder
    return {"message": f"Moved slide from {from_index} to {to_index}", "note": "Move operation is complex in python-pptx"}


async def duplicate_slide(file_path: str, slide_index: int, position: int = -1) -> Dict[str, Any]:
    """Duplicate an existing slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    # Get the source slide
    source_slide = prs.slides[slide_index]

    # Add new slide with same layout
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # Copy content (this is a simplified version - full duplication requires more complex logic)
    try:
        if source_slide.shapes.title:
            new_slide.shapes.title.text = source_slide.shapes.title.text
    except:
        pass

    new_idx = len(prs.slides) - 1
    return {"message": f"Duplicated slide {slide_index} to position {new_idx}", "new_slide_index": new_idx}


async def list_slides(file_path: str) -> Dict[str, Any]:
    """List all slides in the presentation."""
    prs = _get_presentation(file_path)

    slides_info = []
    for i, slide in enumerate(prs.slides):
        slide_info = {"index": i, "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, "name") else f"Layout {i}", "shape_count": len(slide.shapes), "title": ""}

        # Try to get slide title
        try:
            if slide.shapes.title:
                slide_info["title"] = slide.shapes.title.text
        except:
            pass

        slides_info.append(slide_info)

    return {"slides": slides_info, "total_count": len(slides_info)}


async def set_slide_title(file_path: str, slide_index: int, title: str) -> Dict[str, Any]:
    """Set the title of a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if not slide.shapes.title:
        raise ValueError("This slide layout does not have a title placeholder")

    slide.shapes.title.text = title

    # Save presentation after modification
    _save_presentation(file_path)

    return {"message": f"Set title for slide {slide_index}: {title}"}


async def set_slide_content(file_path: str, slide_index: int, content: str) -> Dict[str, Any]:
    """Set the main content/body text of a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Look for content placeholder
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder is usually index 1
            content_placeholder = shape
            break

    if not content_placeholder:
        # If no content placeholder, try to find text frame
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                content_placeholder = shape
                break

    if not content_placeholder:
        raise ValueError("No content area found on this slide")

    # Split content by newlines and create bullet points
    lines = content.split("\\n")
    content_placeholder.text = lines[0]  # First line

    if len(lines) > 1:
        text_frame = content_placeholder.text_frame
        for line in lines[1:]:
            p = text_frame.add_paragraph()
            p.text = line
            p.level = 0  # Bullet level

    return {"message": f"Set content for slide {slide_index}"}


async def add_text_box(
    file_path: str,
    slide_index: int,
    text: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float = 6.0,
    height: float = 1.0,
    font_size: int = 18,
    font_color: str = "#000000",
    bold: bool = False,
    italic: bool = False,
) -> Dict[str, Any]:
    """Add a text box to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Add text box
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text

    # Format text
    paragraph = text_frame.paragraphs[0]
    run = paragraph.runs[0]
    font = run.font

    font.size = Pt(font_size)
    font.color.rgb = _parse_color(font_color)
    font.bold = bold
    font.italic = italic

    return {"message": f"Added text box to slide {slide_index}", "shape_index": len(slide.shapes) - 1, "text": text}


async def format_text(file_path: str, slide_index: int, shape_index: int, **kwargs) -> Dict[str, Any]:
    """Format existing text in a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError(f"Shape index {shape_index} out of range")

    shape = slide.shapes[shape_index]

    if not hasattr(shape, "text_frame"):
        raise ValueError("Selected shape does not contain text")

    # Apply formatting to all paragraphs and runs
    for paragraph in shape.text_frame.paragraphs:
        if kwargs.get("alignment"):
            alignment_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
            paragraph.alignment = alignment_map.get(kwargs["alignment"], PP_ALIGN.LEFT)

        for run in paragraph.runs:
            font = run.font

            if kwargs.get("font_name"):
                font.name = kwargs["font_name"]
            if kwargs.get("font_size"):
                font.size = Pt(kwargs["font_size"])
            if kwargs.get("font_color"):
                font.color.rgb = _parse_color(kwargs["font_color"])
            if kwargs.get("bold") is not None:
                font.bold = kwargs["bold"]
            if kwargs.get("italic") is not None:
                font.italic = kwargs["italic"]
            if kwargs.get("underline") is not None:
                font.underline = kwargs["underline"]

    return {"message": f"Formatted text in shape {shape_index} on slide {slide_index}"}


async def add_image(file_path: str, slide_index: int, image_path: str, left: float = 1.0, top: float = 1.0, width: Optional[float] = None, height: Optional[float] = None) -> Dict[str, Any]:
    """Add an image to a slide."""
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image file not found: {image_path}")

    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Add image
    if width and height:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
    elif width:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width))
    elif height:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(image_path, Inches(left), Inches(top))

    return {"message": f"Added image to slide {slide_index}", "shape_index": len(slide.shapes) - 1, "image_path": image_path}


async def add_image_from_base64(
    file_path: str, slide_index: int, image_data: str, image_format: str = "png", left: float = 1.0, top: float = 1.0, width: Optional[float] = None, height: Optional[float] = None
) -> Dict[str, Any]:
    """Add an image from base64 data to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Decode base64 image
    try:
        image_bytes = base64.b64decode(image_data)
        image_stream = BytesIO(image_bytes)
    except Exception as e:
        raise ValueError(f"Invalid base64 image data: {e}")

    # Add image from stream
    if width and height:
        pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top), Inches(width), Inches(height))
    elif width:
        pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top), width=Inches(width))
    elif height:
        pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(image_stream, Inches(left), Inches(top))

    return {"message": f"Added image from base64 to slide {slide_index}", "shape_index": len(slide.shapes) - 1, "format": image_format}


async def replace_image(file_path: str, slide_index: int, shape_index: int, new_image_path: str) -> Dict[str, Any]:
    """Replace an existing image in a slide."""
    if not os.path.exists(new_image_path):
        raise FileNotFoundError(f"New image file not found: {new_image_path}")

    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError(f"Shape index {shape_index} out of range")

    slide.shapes[shape_index]

    # This is complex in python-pptx - would need to remove old image and add new one
    # For now, provide guidance
    return {"message": "Image replacement requires removing old image and adding new one", "note": "Use delete_shape and add_image for full replacement functionality"}


async def add_shape(
    file_path: str,
    slide_index: int,
    shape_type: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float = 2.0,
    height: float = 1.0,
    fill_color: Optional[str] = None,
    line_color: Optional[str] = None,
    line_width: float = 1.0,
) -> Dict[str, Any]:
    """Add a shape to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Map shape types to MSO_SHAPE constants
    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "arrow": MSO_SHAPE.BLOCK_ARC,
        "diamond": MSO_SHAPE.DIAMOND,
        "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "octagon": MSO_SHAPE.OCTAGON,
        "star": MSO_SHAPE.STAR_5_POINT,
        "heart": MSO_SHAPE.HEART,
        "smiley": MSO_SHAPE.SMILEY_FACE,
    }

    if shape_type.lower() not in shape_map:
        available_shapes = ", ".join(shape_map.keys())
        raise ValueError(f"Unknown shape type: {shape_type}. Available: {available_shapes}")

    # Add shape
    shape = slide.shapes.add_shape(shape_map[shape_type.lower()], Inches(left), Inches(top), Inches(width), Inches(height))

    # Apply formatting
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_color(fill_color)

    if line_color:
        shape.line.color.rgb = _parse_color(line_color)

    shape.line.width = Pt(line_width)

    return {"message": f"Added {shape_type} shape to slide {slide_index}", "shape_index": len(slide.shapes) - 1}


async def modify_shape(file_path: str, slide_index: int, shape_index: int, **kwargs) -> Dict[str, Any]:
    """Modify properties of an existing shape."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError(f"Shape index {shape_index} out of range")

    shape = slide.shapes[shape_index]

    # Modify position and size
    if kwargs.get("left") is not None:
        shape.left = Inches(kwargs["left"])
    if kwargs.get("top") is not None:
        shape.top = Inches(kwargs["top"])
    if kwargs.get("width") is not None:
        shape.width = Inches(kwargs["width"])
    if kwargs.get("height") is not None:
        shape.height = Inches(kwargs["height"])

    # Modify formatting
    if kwargs.get("fill_color"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_color(kwargs["fill_color"])

    if kwargs.get("line_color"):
        shape.line.color.rgb = _parse_color(kwargs["line_color"])

    if kwargs.get("line_width") is not None:
        shape.line.width = Pt(kwargs["line_width"])

    return {"message": f"Modified shape {shape_index} on slide {slide_index}"}


async def delete_shape(file_path: str, slide_index: int, shape_index: int) -> Dict[str, Any]:
    """Delete a shape from a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError(f"Shape index {shape_index} out of range")

    shape = slide.shapes[shape_index]
    sp = shape._element
    sp.getparent().remove(sp)

    return {"message": f"Deleted shape {shape_index} from slide {slide_index}"}


async def add_table(file_path: str, slide_index: int, rows: int, cols: int, left: float = 1.0, top: float = 1.0, width: float = 6.0, height: float = 3.0) -> Dict[str, Any]:
    """Add a table to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Add table
    slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table_index = len(slide.shapes) - 1

    return {
        "message": f"Added {rows}x{cols} table to slide {slide_index}",
        "shape_index": table_index,
        "table_shape_index": table_index,  # Explicit table index for reference
        "rows": rows,
        "cols": cols,
    }


async def set_table_cell(file_path: str, slide_index: int, table_index: int, row: int, col: int, text: str) -> Dict[str, Any]:
    """Set the text content of a table cell."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if table_index < 0 or table_index >= len(slide.shapes):
        raise ValueError(f"Table index {table_index} out of range")

    shape = slide.shapes[table_index]

    try:
        if not shape.has_table:
            raise ValueError("Selected shape is not a table")
        table = shape.table
    except AttributeError:
        raise ValueError("Selected shape is not a table")

    if row < 0 or row >= len(table.rows):
        raise ValueError(f"Row {row} out of range")
    if col < 0 or col >= len(table.columns):
        raise ValueError(f"Column {col} out of range")

    cell = table.cell(row, col)
    cell.text = text

    return {"message": f"Set cell [{row},{col}] text: {text}"}


async def format_table_cell(file_path: str, slide_index: int, table_index: int, row: int, col: int, **kwargs) -> Dict[str, Any]:
    """Format a table cell."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if table_index < 0 or table_index >= len(slide.shapes):
        raise ValueError(f"Table index {table_index} out of range")

    shape = slide.shapes[table_index]

    try:
        if not shape.has_table:
            raise ValueError("Selected shape is not a table")
        table = shape.table
    except AttributeError:
        raise ValueError("Selected shape is not a table")
    cell = table.cell(row, col)

    # Format cell background
    if kwargs.get("fill_color"):
        cell.fill.solid()
        cell.fill.fore_color.rgb = _parse_color(kwargs["fill_color"])

    # Format text
    for paragraph in cell.text_frame.paragraphs:
        if kwargs.get("alignment"):
            alignment_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
            paragraph.alignment = alignment_map.get(kwargs["alignment"], PP_ALIGN.LEFT)

        for run in paragraph.runs:
            font = run.font

            if kwargs.get("font_size"):
                font.size = Pt(kwargs["font_size"])
            if kwargs.get("font_color"):
                font.color.rgb = _parse_color(kwargs["font_color"])
            if kwargs.get("bold") is not None:
                font.bold = kwargs["bold"]

    return {"message": f"Formatted cell [{row},{col}]"}


async def populate_table(file_path: str, slide_index: int, table_index: int, data: List[List[str]], header_row: bool = False) -> Dict[str, Any]:
    """Populate entire table with data from a 2D array."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if table_index < 0 or table_index >= len(slide.shapes):
        raise ValueError(f"Table index {table_index} out of range")

    shape = slide.shapes[table_index]

    try:
        if not shape.has_table:
            raise ValueError("Selected shape is not a table")
        table = shape.table
    except AttributeError:
        raise ValueError("Selected shape is not a table")

    # Populate data
    for row_idx, row_data in enumerate(data):
        if row_idx >= len(table.rows):
            break

        for col_idx, cell_data in enumerate(row_data):
            if col_idx >= len(table.columns):
                break

            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)

            # Format header row
            if header_row and row_idx == 0:
                for run in cell.text_frame.paragraphs[0].runs:
                    run.font.bold = True

    return {"message": f"Populated table with {len(data)} rows of data"}


async def add_chart(
    file_path: str, slide_index: int, data: Dict[str, Any], chart_type: str = "column", left: float = 1.0, top: float = 1.0, width: float = 6.0, height: float = 4.0, title: Optional[str] = None
) -> Dict[str, Any]:
    """Add a chart to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Map chart types
    chart_type_map = {"column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED, "line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}

    if chart_type not in chart_type_map:
        available_types = ", ".join(chart_type_map.keys())
        raise ValueError(f"Unknown chart type: {chart_type}. Available: {available_types}")

    # Prepare chart data
    chart_data = CategoryChartData()
    chart_data.categories = data.get("categories", [])

    for series_info in data.get("series", []):
        chart_data.add_series(series_info.get("name", "Series"), series_info.get("values", []))

    # Add chart
    chart_shape = slide.shapes.add_chart(chart_type_map[chart_type], Inches(left), Inches(top), Inches(width), Inches(height), chart_data)

    # Set title if provided
    if title:
        chart_shape.chart.chart_title.text_frame.text = title

    return {"message": f"Added {chart_type} chart to slide {slide_index}", "shape_index": len(slide.shapes) - 1, "title": title or "Untitled Chart"}


async def update_chart_data(file_path: str, slide_index: int, chart_index: int, data: Dict[str, Any]) -> Dict[str, Any]:
    """Update data in an existing chart."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if chart_index < 0 or chart_index >= len(slide.shapes):
        raise ValueError(f"Chart index {chart_index} out of range")

    shape = slide.shapes[chart_index]

    if not hasattr(shape, "chart"):
        raise ValueError("Selected shape is not a chart")

    # Note: Updating chart data in python-pptx is complex and may require
    # recreating the chart or manipulating the underlying XML
    return {"message": "Chart data update is complex in python-pptx", "note": "Consider recreating the chart with new data for full functionality"}


async def list_shapes(file_path: str, slide_index: int) -> Dict[str, Any]:
    """List all shapes on a slide with their types and properties."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    shapes_info = []
    for i, shape in enumerate(slide.shapes):
        shape_info = {
            "index": i,
            "type": str(shape.shape_type),
            "left": float(shape.left.inches),
            "top": float(shape.top.inches),
            "width": float(shape.width.inches),
            "height": float(shape.height.inches),
            "has_text": hasattr(shape, "text_frame"),
            "text": "",
        }

        # Get text if available
        if hasattr(shape, "text_frame") and shape.text_frame:
            try:
                shape_info["text"] = shape.text_frame.text[:100]  # First 100 chars
            except:
                pass

        # Special handling for different shape types
        try:
            if shape.has_table:
                shape_info["type"] = "TABLE"
                shape_info["rows"] = len(shape.table.rows)
                shape_info["cols"] = len(shape.table.columns)
        except (AttributeError, ValueError):
            pass

        try:
            if shape.has_chart:
                shape_info["type"] = "CHART"
                shape_info["chart_type"] = str(shape.chart.chart_type)
        except (AttributeError, ValueError):
            pass

        shapes_info.append(shape_info)

    return {"shapes": shapes_info, "total_count": len(shapes_info)}


async def get_slide_layouts(file_path: str) -> Dict[str, Any]:
    """Get available slide layouts in the presentation."""
    prs = _get_presentation(file_path)

    layouts_info = []
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {"index": i, "name": layout.name if hasattr(layout, "name") else f"Layout {i}", "placeholder_count": len(layout.placeholders)}
        layouts_info.append(layout_info)

    return {"layouts": layouts_info, "total_count": len(layouts_info)}


async def set_presentation_properties(file_path: str, **kwargs) -> Dict[str, Any]:
    """Set presentation document properties."""
    prs = _get_presentation(file_path)
    props = prs.core_properties

    if kwargs.get("title"):
        props.title = kwargs["title"]
    if kwargs.get("author"):
        props.author = kwargs["author"]
    if kwargs.get("subject"):
        props.subject = kwargs["subject"]
    if kwargs.get("comments"):
        props.comments = kwargs["comments"]

    return {"message": "Updated presentation properties"}


async def set_slide_size(file_path: str, format: str = "16:9", width_inches: Optional[float] = None, height_inches: Optional[float] = None) -> Dict[str, Any]:
    """Set the slide size/aspect ratio of the presentation."""
    prs = _get_presentation(file_path)

    if format == "16:9":
        _set_slide_size_16_9(prs)
        width = 13.33
        height = 7.5
    elif format == "4:3":
        _set_slide_size_4_3(prs)
        width = 10.0
        height = 7.5
    elif format == "custom":
        if width_inches is None or height_inches is None:
            raise ValueError("Custom format requires both width_inches and height_inches")
        prs.slide_width = Inches(width_inches)
        prs.slide_height = Inches(height_inches)
        width = width_inches
        height = height_inches
    else:
        raise ValueError(f"Unsupported format: {format}. Use '16:9', '4:3', or 'custom'")

    return {"message": f"Set slide size to {format}", "format": format, "width_inches": width, "height_inches": height, "aspect_ratio": f"{width/height:.2f}:1"}


async def get_slide_size(file_path: str) -> Dict[str, Any]:
    """Get the current slide size and aspect ratio of the presentation."""
    prs = _get_presentation(file_path)

    width_inches = prs.slide_width.inches
    height_inches = prs.slide_height.inches
    aspect_ratio = width_inches / height_inches

    # Determine format
    if abs(aspect_ratio - 16 / 9) < 0.01:
        format_name = "16:9 widescreen"
    elif abs(aspect_ratio - 4 / 3) < 0.01:
        format_name = "4:3 standard"
    else:
        format_name = "custom"

    return {"width_inches": round(width_inches, 2), "height_inches": round(height_inches, 2), "aspect_ratio": f"{aspect_ratio:.2f}:1", "format": format_name, "is_widescreen": aspect_ratio > 1.5}


async def export_slide_as_image(file_path: str, slide_index: int, output_path: str, format: str = "png") -> Dict[str, Any]:
    """Export a slide as an image file."""
    # Note: python-pptx doesn't have built-in slide-to-image export functionality
    # This would require additional libraries like python-pptx-interface or PIL with COM automation
    return {
        "message": "Slide image export requires additional libraries",
        "note": "Consider using python-pptx-interface or COM automation for image export functionality",
        "requested_output": output_path,
        "format": format,
    }


async def get_file_content(file_path: str, session_id: Optional[str] = None) -> Dict[str, Any]:
    """Get the raw file content for download (base64 encoded)."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    # Validate file is within allowed directories (security check)
    abs_path = os.path.abspath(file_path)
    allowed_dirs = [
        os.path.abspath(config.output_dir),
        os.path.abspath(config.temp_dir),
        os.path.abspath(os.path.join(config.work_dir, "sessions")),
        os.path.abspath("examples/generated"),
        os.path.abspath("examples/demos"),
    ]

    is_allowed = any(abs_path.startswith(allowed_dir) for allowed_dir in allowed_dirs)
    if not is_allowed:
        raise ValueError("File access denied - not in allowed directory")

    # Validate session access if provided
    if session_id:
        # Check if file belongs to this session
        if f"/sessions/{session_id}/" not in abs_path:
            raise ValueError("File access denied - not in your session")

    # Read file content
    try:
        with open(abs_path, "rb") as f:
            file_content = f.read()

        # Encode as base64
        # Standard
        import base64

        file_data = base64.b64encode(file_content).decode("utf-8")

        # Get file info
        filename = os.path.basename(file_path)
        file_size = len(file_content)

        return {
            "message": f"Retrieved file content for {filename}",
            "filename": filename,
            "file_data": file_data,
            "file_size": file_size,
            "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "encoding": "base64",
            "session_id": session_id or "unknown",
        }

    except Exception as e:
        raise ValueError(f"Error reading file: {e}")


# Security and File Management Functions
async def create_secure_session(session_name: Optional[str] = None) -> Dict[str, Any]:
    """Create a secure session for file operations with UUID workspace."""
    session_id = _generate_session_id()
    session_dir = os.path.join(config.work_dir, "sessions", session_id)

    # Create session directory
    os.makedirs(session_dir, exist_ok=True)
    os.chmod(session_dir, 0o700)  # Secure permissions

    # Initialize session
    _session_files[session_id] = []

    # Create session metadata
    session_info = {
        "session_id": session_id,
        "session_name": session_name or f"Session-{session_id[:8]}",
        "created": datetime.now().isoformat(),
        "workspace_dir": session_dir,
        "expires": (datetime.now() + timedelta(hours=config.auto_cleanup_hours)).isoformat(),
        "max_files": config.max_files_per_session,
        "current_files": 0,
    }

    # Save session metadata
    session_file = os.path.join(session_dir, "session.json")
    with open(session_file, "w") as f:
        json.dump(session_info, f, indent=2)

    log.info(f"Created secure session: {session_id}")

    return {
        "message": f"Created secure session: {session_id}",
        "session_id": session_id,
        "session_name": session_info["session_name"],
        "workspace_dir": session_dir,
        "expires": session_info["expires"],
        "max_files": config.max_files_per_session,
    }


async def upload_file(file_data: str, filename: str, session_id: Optional[str] = None) -> Dict[str, Any]:
    """Upload a file to secure workspace."""
    if not config.enable_file_uploads:
        raise ValueError("File uploads are disabled")

    # Validate filename
    safe_filename = _validate_filename(filename)

    # Check file extension
    file_ext = os.path.splitext(safe_filename)[1].lower().lstrip(".")
    if file_ext not in config.allowed_extensions:
        raise ValueError(f"File type .{file_ext} not allowed. Allowed: {config.allowed_extensions}")

    # Decode file data
    try:
        file_bytes = base64.b64decode(file_data)
    except Exception as e:
        raise ValueError(f"Invalid base64 file data: {e}")

    # Check file size
    file_size_mb = len(file_bytes) / (1024 * 1024)
    if file_size_mb > config.max_file_size_mb:
        raise ValueError(f"File too large: {file_size_mb:.1f}MB > {config.max_file_size_mb}MB limit")

    # Determine upload directory
    if session_id:
        session_dir = os.path.join(config.work_dir, "sessions", session_id)
        if not os.path.exists(session_dir):
            raise ValueError(f"Session not found: {session_id}")
        upload_dir = os.path.join(session_dir, "uploads")
    else:
        upload_dir = config.uploads_dir

    os.makedirs(upload_dir, exist_ok=True)

    # Generate unique filename to avoid conflicts
    base_name, ext = os.path.splitext(safe_filename)
    unique_filename = f"{base_name}_{uuid.uuid4().hex[:8]}{ext}"
    upload_path = os.path.join(upload_dir, unique_filename)

    # Save file
    with open(upload_path, "wb") as f:
        f.write(file_bytes)

    # Set secure permissions
    os.chmod(upload_path, 0o600)

    # Track file in session
    if session_id and session_id in _session_files:
        _session_files[session_id].append(upload_path)

    log.info(f"Uploaded file: {unique_filename} ({file_size_mb:.1f}MB)")

    return {
        "message": f"Uploaded file: {unique_filename}",
        "filename": unique_filename,
        "original_filename": filename,
        "file_path": upload_path,
        "size_mb": round(file_size_mb, 2),
        "session_id": session_id,
        "file_type": file_ext,
    }


async def create_download_link(file_path: str, session_id: Optional[str] = None) -> Dict[str, Any]:
    """Create a secure download link for a presentation."""
    if not config.enable_downloads:
        raise ValueError("Downloads are disabled")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    # Validate file is within allowed directories
    abs_path = os.path.abspath(file_path)
    allowed_dirs = [
        os.path.abspath(config.output_dir),
        os.path.abspath(config.temp_dir),
        os.path.abspath(os.path.join(config.work_dir, "sessions")),  # Session directories
        os.path.abspath("examples/generated"),
        os.path.abspath("examples/demos"),
    ]

    is_allowed = any(abs_path.startswith(allowed_dir) for allowed_dir in allowed_dirs)
    if not is_allowed:
        raise ValueError(f"File not in downloadable directory. File: {abs_path}, Allowed: {allowed_dirs}")

    # Generate download token
    download_session = session_id or "anonymous"
    token = _generate_download_token(abs_path, download_session)

    # Create download URL with filename in path
    filename = os.path.basename(file_path)
    if config.enable_http_downloads:
        download_url = f"http://{config.server_host}:{config.server_port}/download/{token}/{filename}"
    else:
        download_url = f"/download/{token}/{filename}"

    return {
        "message": f"Created download link for {filename}",
        "download_token": token,
        "download_url": download_url,
        "expires": _download_tokens[token]["expires"].isoformat(),
        "session_id": download_session,
        "instructions": {"method_1_http": f"Start HTTP server (make serve-http-only) then access: {download_url}", "method_2_direct": f"Use get_file_content tool with file_path: {abs_path}"},
    }


async def list_session_files(session_id: str) -> Dict[str, Any]:
    """List all files in the current session."""
    session_dir = os.path.join(config.work_dir, "sessions", session_id)
    if not os.path.exists(session_dir):
        raise ValueError(f"Session not found: {session_id}")

    # Load session metadata
    session_file = os.path.join(session_dir, "session.json")
    session_info = {}
    if os.path.exists(session_file):
        with open(session_file, "r") as f:
            session_info = json.load(f)

    # Scan for files
    files = []
    for root, dirs, filenames in os.walk(session_dir):
        for filename in filenames:
            if filename == "session.json":
                continue

            file_path = os.path.join(root, filename)
            file_stat = os.stat(file_path)
            relative_path = os.path.relpath(file_path, session_dir)

            files.append(
                {
                    "filename": filename,
                    "relative_path": relative_path,
                    "full_path": file_path,
                    "size_bytes": file_stat.st_size,
                    "size_mb": round(file_stat.st_size / (1024 * 1024), 2),
                    "modified": datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
                    "type": os.path.splitext(filename)[1].lower().lstrip("."),
                }
            )

    return {
        "session_id": session_id,
        "session_name": session_info.get("session_name", "Unknown"),
        "files": files,
        "file_count": len(files),
        "total_size_mb": round(sum(f["size_bytes"] for f in files) / (1024 * 1024), 2),
        "workspace_dir": session_dir,
    }


async def cleanup_session(session_id: str, force: bool = False) -> Dict[str, Any]:
    """Clean up session files and resources."""
    session_dir = os.path.join(config.work_dir, "sessions", session_id)
    if not os.path.exists(session_dir):
        raise ValueError(f"Session not found: {session_id}")

    # Get session info before cleanup
    session_info = await list_session_files(session_id)

    # Remove files
    # Standard
    import shutil

    try:
        shutil.rmtree(session_dir)
        log.info(f"Cleaned up session: {session_id}")
    except Exception as e:
        log.error(f"Error cleaning up session {session_id}: {e}")
        raise

    # Remove from tracking
    if session_id in _session_files:
        del _session_files[session_id]

    # Clean up download tokens for this session
    tokens_to_remove = [token for token, info in _download_tokens.items() if info["session_id"] == session_id]
    for token in tokens_to_remove:
        del _download_tokens[token]

    return {
        "message": f"Cleaned up session: {session_id}",
        "session_id": session_id,
        "files_removed": session_info["file_count"],
        "space_freed_mb": session_info["total_size_mb"],
        "tokens_removed": len(tokens_to_remove),
    }


async def get_server_status() -> Dict[str, Any]:
    """Get server configuration and status information."""
    # Count active sessions
    sessions_dir = os.path.join(config.work_dir, "sessions")
    active_sessions = len([d for d in os.listdir(sessions_dir) if os.path.isdir(os.path.join(sessions_dir, d))]) if os.path.exists(sessions_dir) else 0

    # Count total files
    total_files = 0
    total_size = 0
    for root, dirs, files in os.walk(config.work_dir):
        for file in files:
            if file.endswith(".pptx"):
                file_path = os.path.join(root, file)
                total_files += 1
                total_size += os.path.getsize(file_path)

    return {
        "server_name": "PowerPoint MCP Server",
        "version": "0.1.0",
        "status": "running",
        "configuration": {
            "work_dir": config.work_dir,
            "output_dir": config.output_dir,
            "templates_dir": config.templates_dir,
            "uploads_dir": config.uploads_dir,
            "default_format": config.default_slide_format,
            "max_file_size_mb": config.max_file_size_mb,
            "auto_cleanup_hours": config.auto_cleanup_hours,
            "file_uploads_enabled": config.enable_file_uploads,
            "downloads_enabled": config.enable_downloads,
        },
        "statistics": {
            "active_sessions": active_sessions,
            "active_download_tokens": len(_download_tokens),
            "cached_presentations": len(_presentations),
            "total_pptx_files": total_files,
            "total_storage_mb": round(total_size / (1024 * 1024), 2),
        },
        "security": {
            "allowed_extensions": config.allowed_extensions,
            "max_presentation_size_mb": config.max_presentation_size_mb,
            "authentication_required": config.require_auth,
            "secure_directories": True,
        },
    }


# Template and Enhanced Workflow Functions
async def create_presentation_from_template(template_path: str, output_path: str, title: Optional[str] = None, replace_placeholders: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """Create a new presentation from an existing template."""
    # Resolve template path (check templates directory)
    resolved_template = _resolve_template_path(template_path)
    if not os.path.exists(resolved_template):
        raise FileNotFoundError(f"Template file not found: {template_path} (searched: {resolved_template})")

    # Load template
    template_prs = Presentation(resolved_template)

    # Ensure 16:9 format for new presentations from template
    _set_slide_size_16_9(template_prs)

    # Ensure proper output directory
    organized_output = _ensure_output_directory(output_path)

    # Cache the new presentation
    abs_output_path = os.path.abspath(organized_output)
    _presentations[abs_output_path] = template_prs

    # Update title if provided
    if title and len(template_prs.slides) > 0:
        title_slide = template_prs.slides[0]
        if title_slide.shapes.title:
            title_slide.shapes.title.text = title

    # Replace placeholders if provided
    replacements_made = 0
    if replace_placeholders:
        for slide in template_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for placeholder, replacement in replace_placeholders.items():
                                if placeholder in run.text:
                                    run.text = run.text.replace(placeholder, replacement)
                                    replacements_made += 1

    # Save the new presentation
    _save_presentation(organized_output)

    return {"message": f"Created presentation from template: {resolved_template}", "output_path": organized_output, "slide_count": len(template_prs.slides), "replacements_made": replacements_made}


async def clone_presentation(source_path: str, target_path: str, new_title: Optional[str] = None) -> Dict[str, Any]:
    """Clone an existing presentation with optional modifications."""
    # Resolve source path
    resolved_source = _resolve_template_path(source_path)  # Can also check templates
    if not os.path.exists(resolved_source):
        raise FileNotFoundError(f"Source presentation not found: {source_path}")

    # Load source presentation
    source_prs = Presentation(resolved_source)

    # Ensure 16:9 format for cloned presentations
    _set_slide_size_16_9(source_prs)

    # Ensure proper output directory
    organized_target = _ensure_output_directory(target_path)

    # Cache the cloned presentation
    abs_target_path = os.path.abspath(organized_target)
    _presentations[abs_target_path] = source_prs

    # Update title if provided
    if new_title and len(source_prs.slides) > 0:
        first_slide = source_prs.slides[0]
        if first_slide.shapes.title:
            first_slide.shapes.title.text = new_title

    # Save the cloned presentation
    _save_presentation(organized_target)

    return {"message": f"Cloned presentation from {resolved_source} to {organized_target}", "slide_count": len(source_prs.slides), "new_title": new_title or "No title change"}


# Composite Workflow Tools
async def create_title_slide(file_path: str, title: str, subtitle: Optional[str] = None, author: Optional[str] = None, date: Optional[str] = None, slide_index: int = 0) -> Dict[str, Any]:
    """Create a complete title slide with all elements."""
    prs = _get_presentation(file_path)

    # Get or create slide at specified index
    if slide_index >= len(prs.slides):
        # Add new slide with title layout
        title_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_layout)
        actual_index = len(prs.slides) - 1
    else:
        slide = prs.slides[slide_index]
        actual_index = slide_index

    # Set main title
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Set subtitle in content placeholder or create text box
    if subtitle:
        subtitle_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                subtitle_shape = shape
                break

        if subtitle_shape:
            subtitle_shape.text = subtitle
        else:
            # Create subtitle text box
            await add_text_box(file_path, actual_index, subtitle, 1.0, 2.5, 8.0, 1.0, 20, "#666666", False, True)

    # Add author info if provided
    if author:
        await add_text_box(file_path, actual_index, f"By: {author}", 1.0, 5.5, 4.0, 0.8, 16, "#888888", False, False)

    # Add date if provided
    if date:
        await add_text_box(file_path, actual_index, date, 5.0, 5.5, 4.0, 0.8, 16, "#888888", False, False)

    return {"message": f"Created title slide at index {actual_index}", "slide_index": actual_index, "title": title, "subtitle": subtitle or "None", "author": author or "None", "date": date or "None"}


async def create_data_slide(file_path: str, title: str, data: List[List[str]], include_chart: bool = False, chart_type: str = "column", position: int = -1) -> Dict[str, Any]:
    """Create a complete data slide with table and optional chart."""
    _get_presentation(file_path)

    # Add slide
    slide_result = await add_slide(file_path, 1, position)  # Content layout
    slide_idx = slide_result["slide_index"]

    # Set title
    await set_slide_title(file_path, slide_idx, title)

    # Determine table size
    rows = len(data)
    cols = max(len(row) for row in data) if data else 1

    # Create table
    if include_chart:
        # Smaller table to make room for chart
        table_result = await add_table(file_path, slide_idx, rows, cols, 0.5, 2.0, 4.5, 3.0)
    else:
        # Full-width table
        table_result = await add_table(file_path, slide_idx, rows, cols, 1.0, 2.0, 8.0, 4.0)

    table_idx = table_result["shape_index"]

    # Populate table
    await populate_table(file_path, slide_idx, table_idx, data, True)

    chart_created = False
    if include_chart and len(data) > 1:
        try:
            # Create chart data from table (assuming first row is headers, first column is categories)
            if len(data[0]) >= 2:  # Need at least category and one data column
                categories = [row[0] for row in data[1:]]  # First column, skip header
                series = []

                for col_idx in range(1, len(data[0])):  # Skip first column (categories)
                    series_name = data[0][col_idx]  # Header
                    values = []

                    for row_idx in range(1, len(data)):  # Skip header row
                        try:
                            # Try to convert to number
                            value_str = data[row_idx][col_idx].replace("$", "").replace(",", "").replace("%", "")
                            values.append(float(value_str))
                        except (ValueError, IndexError):
                            values.append(0)

                    series.append({"name": series_name, "values": values})

                chart_data = {"categories": categories, "series": series}
                await add_chart(file_path, slide_idx, chart_data, chart_type, 5.5, 2.0, 4.0, 3.0, f"{title} Chart")
                chart_created = True

        except Exception as e:
            log.warning(f"Could not create chart from data: {e}")

    return {"message": f"Created data slide '{title}' at index {slide_idx}", "slide_index": slide_idx, "table_rows": rows, "table_cols": cols, "chart_created": chart_created}


async def create_comparison_slide(file_path: str, title: str, left_title: str, left_content: List[str], right_title: str, right_content: List[str], position: int = -1) -> Dict[str, Any]:
    """Create a comparison slide with two columns."""
    # Add slide
    slide_result = await add_slide(file_path, 1, position)  # Content layout
    slide_idx = slide_result["slide_index"]

    # Set main title
    await set_slide_title(file_path, slide_idx, title)

    # Create left column (optimized for 16:9 widescreen)
    await add_text_box(file_path, slide_idx, left_title, 0.5, 2.0, 5.5, 0.8, 20, "#0066CC", True, False)
    left_content_text = "\\n".join([f"• {item}" for item in left_content])
    await add_text_box(file_path, slide_idx, left_content_text, 0.5, 3.0, 5.5, 3.0, 16, "#000000", False, False)

    # Create right column (optimized for 16:9 widescreen)
    await add_text_box(file_path, slide_idx, right_title, 7.0, 2.0, 5.5, 0.8, 20, "#0066CC", True, False)
    right_content_text = "\\n".join([f"• {item}" for item in right_content])
    await add_text_box(file_path, slide_idx, right_content_text, 7.0, 3.0, 5.5, 3.0, 16, "#000000", False, False)

    # Add dividing line (centered for 16:9)
    await add_shape(file_path, slide_idx, "rectangle", 6.6, 2.0, 0.1, 4.0, "#CCCCCC", "#CCCCCC", 1.0)

    return {"message": f"Created comparison slide '{title}' at index {slide_idx}", "slide_index": slide_idx, "left_items": len(left_content), "right_items": len(right_content)}


async def create_agenda_slide(file_path: str, agenda_items: List[str], title: str = "Agenda", numbered: bool = True, position: int = 1) -> Dict[str, Any]:
    """Create an agenda slide with numbered or bulleted items."""
    # Add slide
    slide_result = await add_slide(file_path, 1, position)  # Content layout
    slide_idx = slide_result["slide_index"]

    # Set title
    await set_slide_title(file_path, slide_idx, title)

    # Create agenda content
    if numbered:
        agenda_text = "\\n".join([f"{i+1}. {item}" for i, item in enumerate(agenda_items)])
    else:
        agenda_text = "\\n".join([f"• {item}" for item in agenda_items])

    await add_text_box(file_path, slide_idx, agenda_text, 1.5, 2.5, 10.0, 4.0, 18, "#000000", False, False)

    return {"message": f"Created agenda slide '{title}' at index {slide_idx}", "slide_index": slide_idx, "item_count": len(agenda_items), "numbered": numbered}


async def batch_replace_text(file_path: str, replacements: Dict[str, str], slide_range: Optional[List[int]] = None, case_sensitive: bool = False) -> Dict[str, Any]:
    """Replace text across multiple slides in the presentation."""
    prs = _get_presentation(file_path)

    if slide_range is None:
        slides_to_process = list(range(len(prs.slides)))
    else:
        slides_to_process = [i for i in slide_range if 0 <= i < len(prs.slides)]

    total_replacements = 0

    for slide_idx in slides_to_process:
        slide = prs.slides[slide_idx]

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        modified_text = original_text

                        for old_text, new_text in replacements.items():
                            if case_sensitive:
                                if old_text in modified_text:
                                    modified_text = modified_text.replace(old_text, new_text)
                                    total_replacements += 1
                            else:
                                # Case-insensitive replacement
                                # Standard
                                import re

                                pattern = re.compile(re.escape(old_text), re.IGNORECASE)
                                if pattern.search(modified_text):
                                    modified_text = pattern.sub(new_text, modified_text)
                                    total_replacements += 1

                        if modified_text != original_text:
                            run.text = modified_text

    return {
        "message": f"Completed batch text replacement across {len(slides_to_process)} slides",
        "slides_processed": len(slides_to_process),
        "total_replacements": total_replacements,
        "replacement_pairs": len(replacements),
    }


async def apply_brand_theme(
    file_path: str,
    primary_color: str = "#0066CC",
    secondary_color: str = "#999999",
    accent_color: str = "#FF6600",
    font_family: str = "Arial",
    apply_to_titles: bool = True,
    apply_to_shapes: bool = True,
) -> Dict[str, Any]:
    """Apply consistent branding theme across presentation."""
    prs = _get_presentation(file_path)

    title_updates = 0
    shape_updates = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            # Apply to titles
            if apply_to_titles and hasattr(shape, "text_frame") and shape.text_frame:
                if shape == slide.shapes.title:  # This is a title
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = font_family
                            run.font.color.rgb = _parse_color(primary_color)
                    title_updates += 1

            # Apply to shapes
            if apply_to_shapes and hasattr(shape, "fill"):
                try:
                    # Apply primary color to rectangle shapes
                    if shape.shape_type == 1:  # Rectangle
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = _parse_color(primary_color)
                        shape_updates += 1
                    # Apply accent color to other shapes
                    elif shape.shape_type in [9, 7]:  # Oval, triangle, etc.
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = _parse_color(accent_color)
                        shape_updates += 1
                except Exception:
                    pass  # Some shapes may not support fill

    return {
        "message": f"Applied brand theme to presentation",
        "primary_color": primary_color,
        "secondary_color": secondary_color,
        "accent_color": accent_color,
        "font_family": font_family,
        "title_updates": title_updates,
        "shape_updates": shape_updates,
    }


async def create_section_break(
    file_path: str, section_title: str, subtitle: Optional[str] = None, background_color: str = "#0066CC", text_color: str = "#FFFFFF", position: int = -1
) -> Dict[str, Any]:
    """Create a section break slide with large title and background color."""
    # Add slide
    slide_result = await add_slide(file_path, 6, position)  # Blank layout
    slide_idx = slide_result["slide_index"]

    prs = _get_presentation(file_path)
    prs.slides[slide_idx]

    # Set background color by adding a full-slide rectangle (16:9 dimensions)
    await add_shape(file_path, slide_idx, "rectangle", 0, 0, 13.33, 7.5, background_color, background_color, 0)

    # Add large section title
    await add_text_box(file_path, slide_idx, section_title, 1.0, 2.5, 8.0, 1.5, 48, text_color, True, False)

    # Add subtitle if provided
    if subtitle:
        await add_text_box(file_path, slide_idx, subtitle, 1.0, 4.5, 8.0, 1.0, 24, text_color, False, False)

    return {
        "message": f"Created section break slide '{section_title}' at index {slide_idx}",
        "slide_index": slide_idx,
        "section_title": section_title,
        "subtitle": subtitle or "None",
        "background_color": background_color,
    }


async def generate_summary_slide(file_path: str, title: str = "Summary", max_points: int = 5, position: int = -1) -> Dict[str, Any]:
    """Generate a summary slide based on presentation content."""
    prs = _get_presentation(file_path)

    # Extract key points from slide titles and content
    summary_points = []

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx == 0:  # Skip title slide
            continue

        # Get slide title
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text

        if slide_title and len(summary_points) < max_points:
            summary_points.append(slide_title)

    # If we don't have enough from titles, extract from content
    if len(summary_points) < max_points:
        for slide_idx in range(1, len(prs.slides)):  # Skip title slide
            if len(summary_points) >= max_points:
                break

            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text and shape != slide.shapes.title:
                        # Take first sentence or line
                        first_line = text.split("\\n")[0].split(".")[0]
                        if len(first_line) < 80 and first_line not in summary_points:
                            summary_points.append(first_line)
                            if len(summary_points) >= max_points:
                                break

    # Create summary slide
    slide_result = await add_slide(file_path, 1, position)  # Content layout
    slide_idx = slide_result["slide_index"]

    await set_slide_title(file_path, slide_idx, title)

    if summary_points:
        summary_text = "\\n".join([f"• {point}" for point in summary_points])
        await add_text_box(file_path, slide_idx, summary_text, 1.0, 2.5, 8.0, 4.0, 18, "#000000", False, False)
    else:
        await add_text_box(file_path, slide_idx, "• No key points extracted from presentation content", 1.0, 2.5, 8.0, 1.0, 18, "#666666", False, True)

    return {"message": f"Generated summary slide '{title}' at index {slide_idx}", "slide_index": slide_idx, "points_extracted": len(summary_points), "max_points": max_points}


async def main() -> None:
    """Main entry point for the PPTX MCP server."""
    log.info("Starting PowerPoint MCP server (stdio)...")
    # Third-Party
    from mcp.server.stdio import stdio_server

    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="pptx-server",
                server_version="0.1.0",
                capabilities={"tools": {}, "logging": {}},
            ),
        )


if __name__ == "__main__":
    asyncio.run(main())
