#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Enhanced PowerPoint MCP Server Demo Script

This script demonstrates all the enhanced capabilities including templates,
composite workflows, batch operations, and professional slide creation.
"""

# Standard
import asyncio
import os
from pathlib import Path
import sys

# Add src to path for imports
sys.path.insert(0, str(Path(__file__).parent / "src"))

# Third-Party
from pptx_server.server import (
    apply_brand_theme,
    batch_replace_text,
    create_agenda_slide,
    create_comparison_slide,
    create_data_slide,
    create_presentation,
    create_presentation_from_template,
    create_section_break,
    create_title_slide,
    get_presentation_info,
    save_presentation,
)


async def create_enhanced_demo():
    """Create comprehensive demo showcasing all enhanced features."""
    print("🚀 Enhanced PowerPoint MCP Server Demo")
    print("=====================================")

    try:
        # 1. Template System Demo
        print("\n📋 1. TEMPLATE SYSTEM")
        print("-" * 40)

        # Create a corporate template
        print("Creating corporate template...")
        template_file = "examples/templates/corporate_template.pptx"
        await create_presentation(template_file, "{{COMPANY}} {{REPORT_TYPE}}")

        # Add template slides with placeholders
        await create_title_slide(template_file, "{{COMPANY}} {{REPORT_TYPE}}", "{{SUBTITLE}}", "{{DEPARTMENT}}", "{{DATE}}")

        await create_agenda_slide(template_file, ["{{TOPIC_1}}", "{{TOPIC_2}}", "{{TOPIC_3}}", "{{TOPIC_4}}"])

        await save_presentation(template_file)
        print("✅ Created corporate template with placeholders")

        # Use template to create specific presentations
        presentations = [
            {
                "name": "examples/generated/q4_financial_report.pptx",
                "replacements": {
                    "{{COMPANY}}": "TechCorp Industries",
                    "{{REPORT_TYPE}}": "Q4 Financial Report",
                    "{{SUBTITLE}}": "Year-End Performance & Outlook",
                    "{{DEPARTMENT}}": "Finance Department",
                    "{{DATE}}": "December 2024",
                    "{{TOPIC_1}}": "Revenue Growth Analysis",
                    "{{TOPIC_2}}": "Cost Optimization Results",
                    "{{TOPIC_3}}": "Market Share Performance",
                    "{{TOPIC_4}}": "2025 Strategic Priorities",
                },
            },
            {
                "name": "examples/generated/hr_quarterly_update.pptx",
                "replacements": {
                    "{{COMPANY}}": "TechCorp Industries",
                    "{{REPORT_TYPE}}": "HR Quarterly Update",
                    "{{SUBTITLE}}": "Talent & Culture Development",
                    "{{DEPARTMENT}}": "Human Resources",
                    "{{DATE}}": "December 2024",
                    "{{TOPIC_1}}": "Team Growth & Hiring",
                    "{{TOPIC_2}}": "Employee Satisfaction",
                    "{{TOPIC_3}}": "Training & Development",
                    "{{TOPIC_4}}": "Benefits & Compensation",
                },
            },
        ]

        for pres in presentations:
            await create_presentation_from_template(template_file, pres["name"], replace_placeholders=pres["replacements"])
            print(f"✅ Generated: {pres['name']}")

        # 2. Professional Workflow Demo
        print("\n🎯 2. PROFESSIONAL WORKFLOWS")
        print("-" * 40)

        # Create comprehensive business presentation
        showcase_file = "examples/demos/business_showcase.pptx"
        await create_presentation(showcase_file)

        # Professional title slide
        await create_title_slide(showcase_file, "Enterprise Solutions Portfolio", "Innovative Technology for Business Growth", "Solutions Architecture Team", "December 2024")
        print("✅ Created professional title slide")

        # Agenda with strategic topics
        strategic_agenda = [
            "Executive Summary & Vision",
            "Market Analysis & Opportunities",
            "Solution Portfolio Overview",
            "Implementation Roadmap",
            "Investment & ROI Analysis",
            "Partnership Strategy",
            "Risk Assessment & Mitigation",
            "Success Metrics & KPIs",
        ]
        await create_agenda_slide(showcase_file, strategic_agenda, "Strategic Overview")
        print("✅ Created comprehensive agenda")

        # Section breaks for organization
        sections = [
            ("MARKET ANALYSIS", "Understanding Our Competitive Landscape", "#1f4e79"),
            ("SOLUTION PORTFOLIO", "Innovative Products & Services", "#0066cc"),
            ("IMPLEMENTATION", "Roadmap to Success", "#2e8b57"),
            ("FINANCIAL OUTLOOK", "Investment & Returns", "#8b0000"),
        ]

        for section_title, subtitle, color in sections:
            await create_section_break(showcase_file, section_title, subtitle, color)
            print(f"✅ Added section: {section_title}")

        # Data-driven slides with charts
        print("Adding data analysis slides...")

        # Market data
        market_data = [
            ["Market Segment", "2023 Revenue", "2024 Revenue", "Growth %"],
            ["Enterprise", "$12.5M", "$16.2M", "+29.6%"],
            ["SMB", "$8.3M", "$11.7M", "+41.0%"],
            ["Government", "$5.1M", "$6.9M", "+35.3%"],
            ["Healthcare", "$3.2M", "$4.8M", "+50.0%"],
        ]
        await create_data_slide(showcase_file, "Market Segment Performance", market_data, include_chart=True, chart_type="column")

        # Solution comparison
        await create_comparison_slide(
            showcase_file,
            "Current vs Future State",
            "Current Challenges",
            ["Legacy system limitations", "Manual process inefficiencies", "Scattered data sources", "Limited scalability", "High operational costs"],
            "Future Benefits",
            ["Modern, integrated platform", "Automated workflows", "Unified data ecosystem", "Cloud-native scalability", "Optimized cost structure"],
        )
        print("✅ Added comparison analysis")

        # 3. Batch Operations Demo
        print("\n⚡ 3. BATCH OPERATIONS")
        print("-" * 40)

        # Apply consistent terminology
        terminology_updates = {"2024": "FY2024", "Revenue": "Net Revenue", "Growth": "YoY Growth", "Enterprise Solutions": "Enterprise Cloud Solutions", "Implementation": "Deployment"}

        result = await batch_replace_text(showcase_file, terminology_updates)
        print(f"✅ Updated terminology: {result['total_replacements']} changes across {result['slides_processed']} slides")

        # Apply corporate branding
        brand_result = await apply_brand_theme(
            showcase_file, primary_color="#1f4e79", secondary_color="#666666", accent_color="#ff6600", font_family="Calibri"  # Corporate blue  # Professional gray  # Action orange
        )
        print(f"✅ Applied brand theme: {brand_result['title_updates']} titles, {brand_result['shape_updates']} shapes")

        # 4. Save and Generate Reports
        print("\n📊 4. FINAL RESULTS")
        print("-" * 40)

        # Save all presentations
        presentations_to_save = [template_file, "examples/generated/q4_financial_report.pptx", "examples/generated/hr_quarterly_update.pptx", showcase_file]

        results = {}
        for pres_file in presentations_to_save:
            await save_presentation(pres_file)
            if os.path.exists(pres_file):
                info = await get_presentation_info(pres_file)
                size = os.path.getsize(pres_file)
                results[pres_file] = {"slides": info["slide_count"], "size_bytes": size}

        # Summary report
        print("\n🎉 ENHANCED DEMO COMPLETE!")
        print("=" * 50)

        total_slides = sum(r["slides"] for r in results.values())
        total_size = sum(r["size_bytes"] for r in results.values())

        print(f"📊 Generated {len(results)} presentations:")
        print(f"   📄 Total slides: {total_slides}")
        print(f"   💾 Total size: {total_size:,} bytes")
        print()

        for filename, stats in results.items():
            size_kb = stats["size_bytes"] / 1024
            print(f"   • {filename}")
            print(f"     └ {stats['slides']} slides, {size_kb:.1f} KB")

        # Feature showcase summary
        print(f"\n✨ Features Demonstrated:")
        print(f"   🏗️  Template system with placeholder replacement")
        print(f"   🎯 Professional slide workflows (title, agenda, sections)")
        print(f"   📊 Integrated data visualization with charts")
        print(f"   🔄 Batch text replacement across presentations")
        print(f"   🎨 Brand theme application")
        print(f"   📋 Comparison and analysis layouts")
        print(f"   🚀 Enterprise-grade presentation automation")

        # Verification
        print(f"\n🔍 Verification:")
        for filename in results.keys():
            try:
                # Third-Party
                from pptx import Presentation

                prs = Presentation(filename)
                print(f"   ✅ {filename}: Valid ({len(prs.slides)} slides)")
            except Exception as e:
                print(f"   ❌ {filename}: Error - {e}")

        return True

    except Exception as e:
        print(f"\n❌ Demo failed with error: {e}")
        # Standard
        import traceback

        traceback.print_exc()
        return False


async def main():
    """Main demo execution."""
    success = await create_enhanced_demo()

    if success:
        print(f"\n🏆 Enhanced PowerPoint MCP Server Demo completed successfully!")
        print(f"   Open the generated .pptx files to see the professional results.")
        return 0
    else:
        print(f"\n💥 Demo encountered errors.")
        return 1


if __name__ == "__main__":
    exit_code = asyncio.run(main())
    sys.exit(exit_code)
